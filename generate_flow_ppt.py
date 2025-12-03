"""
Generate Move & Groove Lottery Flow Presentation with Visual Mockups
Creates a comprehensive PPT documenting the entire lottery process
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

def create_flow_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    PINK = RGBColor(233, 30, 99)
    PURPLE = RGBColor(156, 39, 176)
    GREEN = RGBColor(76, 175, 80)
    ORANGE = RGBColor(255, 152, 0)
    DARK = RGBColor(33, 33, 33)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    GRAY = RGBColor(158, 158, 158)
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Gradient-like background (using rectangle)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PINK
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(54)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_mockup_slide(title, mockup_elements, description="", highlight_color=PINK):
        """Create a slide with visual mockup of the app screen"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = highlight_color
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(8), Inches(0.6))
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        
        # Mockup area (browser-like frame)
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1), Inches(8), Inches(6))
        frame.fill.solid()
        frame.fill.fore_color.rgb = LIGHT_GRAY
        frame.line.color.rgb = GRAY
        
        # Add mockup elements
        y_pos = 1.3
        for element in mockup_elements:
            elem_type = element.get("type", "text")
            
            if elem_type == "header":
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y_pos), Inches(7.6), Inches(0.6))
                box.fill.solid()
                box.fill.fore_color.rgb = element.get("color", PINK)
                box.line.fill.background()
                
                txt = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(7.6), Inches(0.5))
                tf = txt.text_frame
                tf.paragraphs[0].text = element.get("text", "")
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = WHITE
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                y_pos += 0.8
                
            elif elem_type == "card":
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(element.get("x", 0.7)), Inches(y_pos), Inches(element.get("w", 2.3)), Inches(1.2))
                box.fill.solid()
                box.fill.fore_color.rgb = WHITE
                box.line.color.rgb = element.get("border", PINK)
                
                txt = slide.shapes.add_textbox(Inches(element.get("x", 0.7) + 0.1), Inches(y_pos + 0.2), Inches(element.get("w", 2.3) - 0.2), Inches(0.9))
                tf = txt.text_frame
                tf.paragraphs[0].text = element.get("title", "")
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = element.get("border", PINK)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                if element.get("subtitle"):
                    p = tf.add_paragraph()
                    p.text = element.get("subtitle", "")
                    p.font.size = Pt(10)
                    p.font.color.rgb = DARK
                    p.alignment = PP_ALIGN.CENTER
                
            elif elem_type == "button":
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(element.get("x", 2)), Inches(y_pos), Inches(element.get("w", 4)), Inches(0.5))
                box.fill.solid()
                box.fill.fore_color.rgb = element.get("color", PINK)
                box.line.fill.background()
                
                txt = slide.shapes.add_textbox(Inches(element.get("x", 2)), Inches(y_pos + 0.1), Inches(element.get("w", 4)), Inches(0.4))
                tf = txt.text_frame
                tf.paragraphs[0].text = element.get("text", "Button")
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = WHITE
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                y_pos += 0.7
                
            elif elem_type == "animation_box":
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y_pos), Inches(5.5), Inches(2))
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(26, 26, 46)
                box.line.fill.background()
                
                # Number display
                num_txt = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.3), Inches(5.5), Inches(1))
                tf = num_txt.text_frame
                tf.paragraphs[0].text = element.get("number", "0123")
                tf.paragraphs[0].font.size = Pt(48)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(0, 255, 136)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Counter
                counter_txt = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 1.2), Inches(5.5), Inches(0.4))
                tf = counter_txt.text_frame
                tf.paragraphs[0].text = element.get("counter", "1523 / 2707")
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = GRAY
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Progress bar
                prog_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y_pos + 1.6), Inches(4.5), Inches(0.15))
                prog_bg.fill.solid()
                prog_bg.fill.fore_color.rgb = RGBColor(51, 51, 51)
                prog_bg.line.fill.background()
                
                prog_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y_pos + 1.6), Inches(element.get("progress", 3)), Inches(0.15))
                prog_fill.fill.solid()
                prog_fill.fill.fore_color.rgb = RGBColor(0, 255, 136)
                prog_fill.line.fill.background()
                
                y_pos += 2.3
                
            elif elem_type == "winner_card":
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y_pos), Inches(5.5), Inches(1.2))
                box.fill.solid()
                box.fill.fore_color.rgb = GREEN
                box.line.fill.background()
                
                txt = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.15), Inches(5.5), Inches(1))
                tf = txt.text_frame
                tf.paragraphs[0].text = element.get("title", "PEMENANG #1")
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = WHITE
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                p = tf.add_paragraph()
                p.text = element.get("number", "03997")
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
                
                p2 = tf.add_paragraph()
                p2.text = element.get("details", "Budi Santoso | ****7890")
                p2.font.size = Pt(12)
                p2.font.color.rgb = WHITE
                p2.alignment = PP_ALIGN.CENTER
                
                y_pos += 1.4
                
            elif elem_type == "newrow":
                y_pos += element.get("height", 1.4)
                
            elif elem_type == "text":
                txt = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(7.6), Inches(0.5))
                tf = txt.text_frame
                tf.paragraphs[0].text = element.get("text", "")
                tf.paragraphs[0].font.size = Pt(element.get("size", 14))
                tf.paragraphs[0].font.color.rgb = element.get("color", DARK)
                tf.paragraphs[0].alignment = element.get("align", PP_ALIGN.LEFT)
                y_pos += 0.5
        
        # Description panel on the right
        if description:
            desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1), Inches(4.4), Inches(6))
            desc_box.fill.solid()
            desc_box.fill.fore_color.rgb = WHITE
            desc_box.line.color.rgb = GRAY
            
            desc_title = slide.shapes.add_textbox(Inches(8.9), Inches(1.2), Inches(4), Inches(0.5))
            tf = desc_title.text_frame
            tf.paragraphs[0].text = "Keterangan:"
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = highlight_color
            
            desc_text = slide.shapes.add_textbox(Inches(8.9), Inches(1.7), Inches(4), Inches(5))
            tf = desc_text.text_frame
            tf.word_wrap = True
            
            for i, line in enumerate(description.split("\n")):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(13)
                p.font.color.rgb = DARK
                p.space_after = Pt(6)
        
        return slide
    
    # ========== SLIDES ==========
    
    # 1. Cover slide
    add_title_slide(
        "SISTEM UNDIAN MOVE & GROOVE",
        "7 Desember 2024 | Panduan Alur Pengundian"
    )
    
    # 2. Overview slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "OVERVIEW SISTEM UNDIAN"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # 3 mode boxes
    modes = [
        {"title": "E-VOUCHER", "count": "700", "desc": "4 kategori\n175 per kategori", "color": GREEN},
        {"title": "SHUFFLE", "count": "90", "desc": "3 sesi\n30 per sesi", "color": ORANGE},
        {"title": "WHEEL", "count": "10", "desc": "Grand Prize\nTermurah → Termahal", "color": PURPLE}
    ]
    
    for i, mode in enumerate(modes):
        x = 0.8 + i * 4.2
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.8), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = mode["color"]
        box.line.fill.background()
        
        txt = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(3.8), Inches(2.3))
        tf = txt.text_frame
        tf.paragraphs[0].text = mode["title"]
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = mode["count"]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = mode["desc"]
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
    
    # Total
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(4.3), Inches(5.33), Inches(1))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = PINK
    total_box.line.fill.background()
    
    total_txt = slide.shapes.add_textbox(Inches(4), Inches(4.5), Inches(5.33), Inches(0.7))
    tf = total_txt.text_frame
    tf.paragraphs[0].text = "TOTAL: 800 PEMENANG"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Features
    features_txt = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
    tf = features_txt.text_frame
    tf.paragraphs[0].text = "Fitur Utama:"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    
    features = ["Auto-save hasil undian", "Backup ke Google Drive", "Validasi duplikat", "Transparansi animasi", "Export Excel & PPT"]
    for feat in features:
        p = tf.add_paragraph()
        p.text = "✓ " + feat
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
    
    # 3. Upload Data Screen
    add_mockup_slide(
        "LANGKAH 1: UPLOAD DATA PESERTA",
        [
            {"type": "header", "text": "SISTEM UNDIAN MOVE & GROOVE", "color": PINK},
            {"type": "text", "text": "Upload data peserta undian:", "size": 16},
            {"type": "button", "text": "📁 Upload CSV File", "x": 1, "w": 3, "color": GRAY},
            {"type": "text", "text": "atau masukkan URL Google Sheets:", "size": 14},
            {"type": "button", "text": "🔗 https://docs.google.com/spreadsheets/...", "x": 1, "w": 6.5, "color": RGBColor(66, 133, 244)},
            {"type": "newrow", "height": 0.3},
            {"type": "button", "text": "✅ LOAD DATA", "x": 2.5, "w": 3.5, "color": GREEN},
        ],
        "1. Pilih salah satu:\n   - Upload file CSV\n   - Masukkan URL Google Sheets\n\n2. Format data:\n   - Nomor Undian\n   - Nama\n   - No HP\n\n3. VIP/F otomatis terfilter\n\n4. Klik LOAD DATA\n\n5. Sistem menampilkan:\n   - Total peserta\n   - Peserta eligible\n   - Peserta VIP/F",
        PINK
    )
    
    # 4. Home Screen with 3 panels
    add_mockup_slide(
        "LANGKAH 2: HALAMAN UTAMA",
        [
            {"type": "header", "text": "✅ 3500 peserta (3407 eligible, 93 VIP/F)", "color": GREEN},
            {"type": "card", "title": "E-VOUCHER", "subtitle": "700 Hadiah\n4 Kategori", "x": 0.7, "w": 2.3, "border": GREEN},
            {"type": "card", "title": "SHUFFLE", "subtitle": "90 Hadiah\n3 Sesi", "x": 3.2, "w": 2.3, "border": ORANGE},
            {"type": "card", "title": "WHEEL", "subtitle": "10 Grand Prize", "x": 5.7, "w": 2.3, "border": PURPLE},
        ],
        "Halaman utama menampilkan:\n\n1. Status data peserta\n   (total, eligible, VIP/F)\n\n2. Tiga panel mode undian:\n   - E-VOUCHER (hijau)\n   - SHUFFLE (orange)\n   - WHEEL (ungu)\n\n3. E-Voucher selalu aktif\n\n4. Shuffle & Wheel aktif\n   setelah E-Voucher selesai\n\n5. Tombol hasil untuk\n   melihat pemenang",
        PINK
    )
    
    # 5. E-Voucher Preview
    add_mockup_slide(
        "LANGKAH 3: E-VOUCHER - PREVIEW",
        [
            {"type": "header", "text": "🎁 UNDIAN E-VOUCHER - 700 HADIAH", "color": GREEN},
            {"type": "card", "title": "TOKOPEDIA", "subtitle": "Rp.100.000\n175 pemenang", "x": 0.7, "w": 1.8, "border": GREEN},
            {"type": "card", "title": "INDOMARET", "subtitle": "Rp.100.000\n175 pemenang", "x": 2.6, "w": 1.8, "border": GREEN},
            {"type": "card", "title": "BENSIN", "subtitle": "Rp.100.000\n175 pemenang", "x": 4.5, "w": 1.8, "border": GREEN},
            {"type": "card", "title": "SNL", "subtitle": "Rp.100.000\n175 pemenang", "x": 6.4, "w": 1.8, "border": GREEN},
            {"type": "newrow", "height": 1.5},
            {"type": "button", "text": "🎲 MULAI UNDIAN E-VOUCHER", "x": 1.5, "w": 5.5, "color": GREEN},
        ],
        "Preview E-Voucher:\n\n1. 4 kategori hadiah:\n   - Tokopedia Rp.100.000\n   - Indomaret Rp.100.000\n   - Bensin Rp.100.000\n   - SNL Rp.100.000\n\n2. Masing-masing 175 pemenang\n\n3. Total 700 pemenang\n\n4. Klik tombol hijau\n   untuk mulai undian\n\n5. Animasi cascade akan\n   menampilkan pemenang",
        GREEN
    )
    
    # 6. E-Voucher Animation
    add_mockup_slide(
        "LANGKAH 3: E-VOUCHER - ANIMASI",
        [
            {"type": "header", "text": "🎁 MENGUNDI 700 PEMENANG...", "color": GREEN},
            {"type": "text", "text": "Animasi cascade menampilkan semua pemenang:", "size": 14},
            {"type": "newrow", "height": 0.2},
            {"type": "card", "title": "#1: 0234", "subtitle": "Ahmad H.\nTokopedia", "x": 0.7, "w": 1.5, "border": GREEN},
            {"type": "card", "title": "#2: 1567", "subtitle": "Budi S.\nTokopedia", "x": 2.3, "w": 1.5, "border": GREEN},
            {"type": "card", "title": "#3: 0891", "subtitle": "Citra D.\nIndomaret", "x": 3.9, "w": 1.5, "border": GREEN},
            {"type": "card", "title": "#4: 2345", "subtitle": "Dian P.\nIndomaret", "x": 5.5, "w": 1.5, "border": GREEN},
            {"type": "card", "title": "#5: 0123", "subtitle": "Eko P.\nBensin", "x": 7.1, "w": 1.5, "border": GREEN},
        ],
        "Proses Undian E-Voucher:\n\n1. Progress bar menunjukkan\n   proses shuffle\n\n2. Animasi cascade\n   menampilkan pemenang\n   satu per satu\n\n3. Setiap kartu berisi:\n   - Nomor urut\n   - Nomor undian\n   - Nama peserta\n   - Kategori hadiah\n\n4. 700 pemenang ditampilkan\n\n5. Setelah selesai:\n   Download Excel/PPT",
        GREEN
    )
    
    # 7. Shuffle Mode
    add_mockup_slide(
        "LANGKAH 4: SHUFFLE - 3 SESI",
        [
            {"type": "header", "text": "🔀 UNDIAN SHUFFLE - 90 HADIAH", "color": ORANGE},
            {"type": "text", "text": "Sesi: 0/3 | Sisa: 2707 peserta", "size": 14, "color": GRAY},
            {"type": "newrow", "height": 0.2},
            {"type": "card", "title": "SESI 1", "subtitle": "30 Hadiah\n[Input nama hadiah]", "x": 0.7, "w": 2.3, "border": ORANGE},
            {"type": "card", "title": "SESI 2", "subtitle": "30 Hadiah\nMenunggu Sesi 1", "x": 3.2, "w": 2.3, "border": GRAY},
            {"type": "card", "title": "SESI 3", "subtitle": "30 Hadiah\nMenunggu Sesi 2", "x": 5.7, "w": 2.3, "border": GRAY},
            {"type": "newrow", "height": 1.5},
            {"type": "button", "text": "🎰 PUTAR UNDIAN SESI 1", "x": 1.5, "w": 5.5, "color": ORANGE},
        ],
        "Mode Shuffle:\n\n1. 3 sesi pengundian\n   (masing-masing 30 hadiah)\n\n2. Setiap sesi:\n   - Input nama hadiah\n   - Klik PUTAR UNDIAN\n   - Animasi slot-machine\n   - 30 pemenang muncul\n\n3. Sesi berurutan\n   (harus selesai Sesi 1\n   untuk lanjut Sesi 2)\n\n4. Download per sesi\n   atau gabungan\n\n5. Total 90 pemenang",
        ORANGE
    )
    
    # 8. Shuffle Animation
    add_mockup_slide(
        "LANGKAH 4: SHUFFLE - ANIMASI",
        [
            {"type": "header", "text": "🎰 SESI 1: DOORPRIZE SPESIAL", "color": ORANGE},
            {"type": "animation_box", "number": "2847", "counter": "1523 / 2707", "progress": 2.5},
            {"type": "text", "text": "Animasi menampilkan semua nomor peserta sisa secara berurutan", "size": 12, "color": GRAY},
        ],
        "Animasi Shuffle:\n\n1. Badge menunjukkan\n   total peserta sisa\n\n2. Nomor ditampilkan\n   secara BERURUTAN\n   (bukan random sampling)\n\n3. Counter real-time:\n   '1523 / 2707'\n\n4. Progress bar mengikuti\n   nomor yang sudah lewat\n\n5. Kecepatan:\n   - Cepat di awal\n   - Melambat di akhir\n\n6. Berhenti di pemenang\n   dengan warna emas",
        ORANGE
    )
    
    # 9. Wheel Mode
    add_mockup_slide(
        "LANGKAH 5: WHEEL - 10 GRAND PRIZE",
        [
            {"type": "header", "text": "🎡 HADIAH UTAMA - 10 GRAND PRIZE", "color": PURPLE},
            {"type": "text", "text": "Progress: ▶1  2  3  4  5  6  7  8  9  10", "size": 14, "color": GRAY},
            {"type": "newrow", "height": 0.2},
            {"type": "card", "title": "#1", "subtitle": "Blender\nRp.300.000", "x": 0.7, "w": 1.5, "border": PURPLE},
            {"type": "card", "title": "#2", "subtitle": "Rice Cooker\nRp.500.000", "x": 2.3, "w": 1.5, "border": PURPLE},
            {"type": "card", "title": "...", "subtitle": "", "x": 3.9, "w": 1, "border": GRAY},
            {"type": "card", "title": "#10", "subtitle": "LED TV 43\"\nRp.5.000.000", "x": 5, "w": 1.8, "border": PINK},
            {"type": "newrow", "height": 1.5},
            {"type": "button", "text": "🎡 PUTAR UNDIAN!", "x": 1.5, "w": 5.5, "color": PURPLE},
        ],
        "Mode Wheel:\n\n1. 10 Grand Prize\n   (termurah → termahal)\n\n2. Konfigurasi hadiah:\n   - No\n   - Nama Hadiah\n   - Keterangan\n\n3. Urutan penting:\n   LED TV sebagai\n   hadiah terakhir (klimaks)\n\n4. Setiap spin:\n   - Animasi penuh\n   - Semua nomor diundi\n   - Transparansi penuh\n\n5. Progress bar 1-10",
        PURPLE
    )
    
    # 10. Wheel Animation
    add_mockup_slide(
        "LANGKAH 5: WHEEL - ANIMASI TRANSPARAN",
        [
            {"type": "header", "text": "Mengundi dari 2617 peserta", "color": PURPLE},
            {"type": "animation_box", "number": "03997", "counter": "2617 / 2617", "progress": 4.5},
            {"type": "winner_card", "title": "🎉 PEMENANG #1", "number": "03997", "details": "Budi Santoso | ****7890"},
        ],
        "Transparansi Wheel:\n\n1. Badge 'Mengundi dari\n   X peserta' menunjukkan\n   total pool sisa\n\n2. Animasi sequential:\n   Semua nomor ditampilkan\n   satu per satu\n\n3. Counter: '1/2617',\n   '2/2617', dst\n\n4. Penonton melihat\n   SETIAP nomor diundi\n\n5. Pemenang ditampilkan:\n   - Nomor undian\n   - Nama lengkap\n   - No HP (masked)",
        PURPLE
    )
    
    # 11. Validation & Download
    add_mockup_slide(
        "LANGKAH 6: VALIDASI & DOWNLOAD",
        [
            {"type": "header", "text": "🏆 SELESAI - 10/10 HADIAH TERPILIH", "color": GREEN},
            {"type": "text", "text": "Validasi dan download hasil undian:", "size": 16},
            {"type": "newrow", "height": 0.3},
            {"type": "button", "text": "✅ VALIDASI DUPLIKAT", "x": 0.7, "w": 2.5, "color": GREEN},
            {"type": "button", "text": "📊 EXCEL LENGKAP", "x": 3.3, "w": 2.5, "color": RGBColor(33, 150, 83)},
            {"type": "button", "text": "📽️ PPT LENGKAP", "x": 5.9, "w": 2.5, "color": ORANGE},
            {"type": "newrow", "height": 0.5},
            {"type": "text", "text": "✅ Validasi: Tidak ada duplikat di 800 pemenang", "size": 14, "color": GREEN},
        ],
        "Validasi & Download:\n\n1. VALIDASI DUPLIKAT\n   - Cek semua tahap\n   - Pastikan tidak ada\n     nomor menang 2x\n\n2. EXCEL LENGKAP\n   - Semua pemenang\n   - Sheet per tahap\n   - Detail lengkap\n\n3. PPT LENGKAP\n   - Slide per tahap\n   - Siap presentasi\n\n4. Auto-save aktif\n   - Backup lokal\n   - Backup Google Drive",
        GREEN
    )
    
    # 12. Closing
    add_title_slide(
        "MOVE & GROOVE 2024",
        "Selamat Mengundi! 🎉"
    )
    
    return prs

if __name__ == "__main__":
    prs = create_flow_presentation()
    prs.save("MoveGroove_Flow_Presentation.pptx")
    print("Presentation saved: MoveGroove_Flow_Presentation.pptx")
