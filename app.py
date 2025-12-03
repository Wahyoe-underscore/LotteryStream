import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import secrets
import hashlib
from io import BytesIO, StringIO
import time
import re
import requests
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PRIZE_CONFIG_FILE = "prize_config.json"
LOTTERY_RESULTS_DIR = "lottery_backups"
GDRIVE_FOLDER_NAME = "Move&Groove_Lottery_Results"

# Permanent Google Sheets URL for Move & Groove Dec 7th Event
DEFAULT_SHEETS_URL = "https://docs.google.com/spreadsheets/d/1blM4h0mr4jG2rsphJFs5kqC2rKPO5tl0m4A8SKNcB7E/edit?gid=1638013732#gid=1638013732"

def get_google_drive_access_token():
    """Get access token for Google Drive API"""
    hostname = os.environ.get("REPLIT_CONNECTORS_HOSTNAME")
    x_replit_token = None
    
    if os.environ.get("REPL_IDENTITY"):
        x_replit_token = "repl " + os.environ.get("REPL_IDENTITY")
    elif os.environ.get("WEB_REPL_RENEWAL"):
        x_replit_token = "depl " + os.environ.get("WEB_REPL_RENEWAL")
    
    if not x_replit_token or not hostname:
        return None
    
    try:
        response = requests.get(
            f"https://{hostname}/api/v2/connection?include_secrets=true&connector_names=google-drive",
            headers={
                "Accept": "application/json",
                "X_REPLIT_TOKEN": x_replit_token
            }
        )
        data = response.json()
        connection = data.get("items", [{}])[0] if data.get("items") else {}
        settings = connection.get("settings", {})
        
        access_token = settings.get("access_token") or settings.get("oauth", {}).get("credentials", {}).get("access_token")
        return access_token
    except Exception as e:
        return None

def get_or_create_gdrive_folder(access_token):
    """Get or create the lottery results folder in Google Drive"""
    headers = {"Authorization": f"Bearer {access_token}"}
    
    # Search for existing folder
    search_url = "https://www.googleapis.com/drive/v3/files"
    params = {
        "q": f"name='{GDRIVE_FOLDER_NAME}' and mimeType='application/vnd.google-apps.folder' and trashed=false",
        "spaces": "drive"
    }
    
    try:
        response = requests.get(search_url, headers=headers, params=params)
        files = response.json().get("files", [])
        
        if files:
            return files[0]["id"]
        
        # Create new folder
        create_url = "https://www.googleapis.com/drive/v3/files"
        folder_metadata = {
            "name": GDRIVE_FOLDER_NAME,
            "mimeType": "application/vnd.google-apps.folder"
        }
        response = requests.post(create_url, headers=headers, json=folder_metadata)
        return response.json().get("id")
    except Exception as e:
        return None

def save_to_google_drive(filename, content, access_token, folder_id=None):
    """Upload a file to Google Drive"""
    headers = {"Authorization": f"Bearer {access_token}"}
    
    # Check if file exists
    search_url = "https://www.googleapis.com/drive/v3/files"
    q = f"name='{filename}' and trashed=false"
    if folder_id:
        q += f" and '{folder_id}' in parents"
    
    try:
        response = requests.get(search_url, headers=headers, params={"q": q})
        files = response.json().get("files", [])
        
        if files:
            # Update existing file
            file_id = files[0]["id"]
            upload_url = f"https://www.googleapis.com/upload/drive/v3/files/{file_id}?uploadType=media"
            response = requests.patch(upload_url, headers={**headers, "Content-Type": "application/json"}, data=content)
        else:
            # Create new file
            metadata = {"name": filename}
            if folder_id:
                metadata["parents"] = [folder_id]
            
            # Multipart upload
            boundary = "----WebKitFormBoundary7MA4YWxkTrZu0gW"
            body = (
                f"--{boundary}\r\n"
                f'Content-Type: application/json; charset=UTF-8\r\n\r\n'
                f'{json.dumps(metadata)}\r\n'
                f"--{boundary}\r\n"
                f"Content-Type: application/json\r\n\r\n"
                f"{content}\r\n"
                f"--{boundary}--"
            )
            
            upload_url = "https://www.googleapis.com/upload/drive/v3/files?uploadType=multipart"
            response = requests.post(
                upload_url,
                headers={**headers, "Content-Type": f"multipart/related; boundary={boundary}"},
                data=body.encode()
            )
        
        return response.status_code in [200, 201]
    except Exception as e:
        return False

PRIZE_TIERS = [
    {"name": "Tokopedia Rp.100.000,-", "icon": "🛒", "count": 175, "start": 1, "end": 175},
    {"name": "Indomaret Rp.100.000,-", "icon": "🏪", "count": 175, "start": 176, "end": 350},
    {"name": "Bensin Rp.100.000,-", "icon": "⛽", "count": 175, "start": 351, "end": 525},
    {"name": "SNL Rp.100.000,-", "icon": "🎵", "count": 175, "start": 526, "end": 700},
]

SHUFFLE_CONFIG = [
    {"name": "Sesi 1", "count": 30, "prize": ""},
    {"name": "Sesi 2", "count": 30, "prize": ""},
    {"name": "Sesi 3", "count": 30, "prize": ""},
]

WHEEL_CONFIG = {"count": 10}

def load_prize_config():
    if os.path.exists(PRIZE_CONFIG_FILE):
        try:
            with open(PRIZE_CONFIG_FILE, 'r') as f:
                return json.load(f)
        except:
            return None
    return None

def save_prize_config(config):
    with open(PRIZE_CONFIG_FILE, 'w') as f:
        json.dump(config, f, indent=2)

def get_current_results_file():
    """Get the current session's results file path with timestamp"""
    if "current_results_file" not in st.session_state:
        # Create new timestamped filename
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        st.session_state["current_results_file"] = f"lottery_{timestamp}.json"
    
    # Ensure directory exists
    if not os.path.exists(LOTTERY_RESULTS_DIR):
        os.makedirs(LOTTERY_RESULTS_DIR)
    
    return os.path.join(LOTTERY_RESULTS_DIR, st.session_state["current_results_file"])

def save_lottery_results():
    """Auto-save all lottery results to JSON file with timestamp and Google Drive"""
    results = {
        "evoucher_done": st.session_state.get("evoucher_done", False),
        "shuffle_done": st.session_state.get("shuffle_done", False),
        "wheel_done": st.session_state.get("wheel_done", False),
        "evoucher_results": None,
        "shuffle_results": st.session_state.get("shuffle_results", {}),
        "wheel_winners": st.session_state.get("wheel_winners", []),
        "wheel_prizes": st.session_state.get("wheel_prizes", []),
        "wheel_config": st.session_state.get("wheel_config", []),
        "remaining_pool": None,
        "participant_data": None,
        "data_source_hash": st.session_state.get("data_source_hash", ""),
        "saved_at": time.strftime("%Y-%m-%d %H:%M:%S"),
    }
    
    # Convert DataFrames to JSON-serializable format
    if "evoucher_results" in st.session_state and st.session_state["evoucher_results"] is not None:
        results["evoucher_results"] = st.session_state["evoucher_results"].to_dict('records')
    
    if "remaining_pool" in st.session_state and st.session_state["remaining_pool"] is not None:
        results["remaining_pool"] = st.session_state["remaining_pool"].to_dict('records')
    
    if "participant_data" in st.session_state and st.session_state["participant_data"] is not None:
        results["participant_data"] = st.session_state["participant_data"].to_dict('records')
    
    results_json = json.dumps(results, indent=2, default=str)
    
    # Save to local file
    results_file = get_current_results_file()
    temp_file = results_file + ".tmp"
    local_saved = False
    try:
        with open(temp_file, 'w') as f:
            f.write(results_json)
        os.replace(temp_file, results_file)
        local_saved = True
    except Exception as e:
        pass
    
    # Save to Google Drive
    gdrive_saved = False
    try:
        access_token = get_google_drive_access_token()
        if access_token:
            folder_id = get_or_create_gdrive_folder(access_token)
            filename = st.session_state.get("current_results_file", "lottery_results.json")
            gdrive_saved = save_to_google_drive(filename, results_json, access_token, folder_id)
            st.session_state["gdrive_save_status"] = gdrive_saved
    except Exception as e:
        st.session_state["gdrive_save_status"] = False
    
    return local_saved or gdrive_saved

def get_latest_results_file():
    """Find the most recent lottery results file"""
    if not os.path.exists(LOTTERY_RESULTS_DIR):
        return None
    
    files = [f for f in os.listdir(LOTTERY_RESULTS_DIR) if f.startswith("lottery_") and f.endswith(".json")]
    if not files:
        return None
    
    # Sort by filename (which includes timestamp)
    files.sort(reverse=True)
    return os.path.join(LOTTERY_RESULTS_DIR, files[0])

def load_lottery_results():
    """Load lottery results from the most recent JSON file"""
    results_file = get_latest_results_file()
    if not results_file or not os.path.exists(results_file):
        return False
    
    try:
        with open(results_file, 'r') as f:
            results = json.load(f)
        
        st.session_state["evoucher_done"] = results.get("evoucher_done", False)
        st.session_state["shuffle_done"] = results.get("shuffle_done", False)
        st.session_state["wheel_done"] = results.get("wheel_done", False)
        st.session_state["shuffle_results"] = results.get("shuffle_results", {})
        st.session_state["wheel_winners"] = results.get("wheel_winners", [])
        st.session_state["wheel_prizes"] = results.get("wheel_prizes", [])
        st.session_state["wheel_config"] = results.get("wheel_config", [])
        st.session_state["data_source_hash"] = results.get("data_source_hash", "")
        
        # Set the current file to the loaded one (to continue saving to same file)
        st.session_state["current_results_file"] = os.path.basename(results_file)
        
        # Restore DataFrames
        if results.get("evoucher_results"):
            st.session_state["evoucher_results"] = pd.DataFrame(results["evoucher_results"])
        
        if results.get("remaining_pool"):
            st.session_state["remaining_pool"] = pd.DataFrame(results["remaining_pool"])
        
        if results.get("participant_data"):
            st.session_state["participant_data"] = pd.DataFrame(results["participant_data"])
        
        return True
    except Exception as e:
        return False

def reset_lottery_session():
    """Reset all lottery data and start a new session with new timestamp"""
    # Clear session state
    keys_to_clear = [
        "evoucher_done", "evoucher_results", 
        "shuffle_done", "shuffle_results",
        "wheel_done", "wheel_winners", "wheel_prizes", "wheel_config",
        "remaining_pool", "participant_data",
        "current_results_file", "results_loaded",
        "data_source_hash", "last_content_hash",
        "sheets_df", "last_sheets_hash"
    ]
    for key in keys_to_clear:
        if key in st.session_state:
            del st.session_state[key]
    
    # Create new timestamp for new session
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    st.session_state["current_results_file"] = f"lottery_{timestamp}.json"
    st.session_state["results_loaded"] = True

def calculate_total_winners(prize_tiers):
    return sum(tier["count"] for tier in prize_tiers)

def secure_shuffle(items):
    items = list(items)
    for i in range(len(items) - 1, 0, -1):
        j = secrets.randbelow(i + 1)
        items[i], items[j] = items[j], items[i]
    return items

def is_eligible_for_prize(name, phone):
    name_str = str(name).strip().upper() if pd.notna(name) else ""
    phone_str = str(phone).strip().upper() if pd.notna(phone) else ""
    
    # Exclude if name contains VIP
    if "VIP" in name_str:
        return False
    
    # Exclude if name is exactly "F" or contains F as a word
    if name_str == "F" or name_str.startswith("F ") or name_str.endswith(" F") or " F " in name_str:
        return False
    
    # Exclude if phone is "F", empty when name is "F", or starts with "F"
    if phone_str == "F" or phone_str.startswith("F"):
        return False
    
    return True

def get_prize_dynamic(rank, prize_tiers):
    for tier in prize_tiers:
        if tier["start"] <= rank <= tier["end"]:
            return tier["name"]
    return "Hadiah"

def format_phone(phone):
    """Format phone number for display (no masking for internal use)"""
    phone_str = str(phone) if pd.notna(phone) else ""
    if phone_str.lower() == "nan":
        return "-"
    return phone_str if phone_str else "-"

def create_shuffle_animation_html(all_participants, winners, prize_name="Hadiah"):
    """Create an animated shuffle display showing winners being selected - cascade style for many winners"""
    winners_js = json.dumps(winners)
    all_nums_js = json.dumps(all_participants[:100])  # Use 100 random numbers for animation variety
    total_winners = len(winners)
    
    # Calculate grid layout based on number of winners
    if total_winners <= 10:
        cols = 5
        slot_size = "90px"
        font_size = "1.3rem"
    elif total_winners <= 20:
        cols = 5
        slot_size = "80px"
        font_size = "1.1rem"
    else:
        cols = 6
        slot_size = "70px"
        font_size = "1rem"
    
    html = f'''
    <!DOCTYPE html>
    <html>
    <head>
        <style>
            * {{ margin: 0; padding: 0; box-sizing: border-box; }}
            body {{
                display: flex;
                flex-direction: column;
                align-items: center;
                justify-content: flex-start;
                min-height: 100vh;
                background: transparent;
                font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                padding: 10px;
            }}
            .container {{
                text-align: center;
                width: 100%;
                max-width: 600px;
            }}
            .header {{
                background: linear-gradient(135deg, #FF9800, #FF5722);
                padding: 15px;
                border-radius: 15px;
                margin-bottom: 15px;
            }}
            .prize-title {{
                font-size: 1.5rem;
                font-weight: 800;
                color: white;
                text-shadow: 2px 2px 4px rgba(0,0,0,0.2);
            }}
            .counter {{
                font-size: 1rem;
                color: #fff;
                margin-top: 5px;
            }}
            .slots-container {{
                display: grid;
                grid-template-columns: repeat({cols}, 1fr);
                gap: 8px;
                margin-bottom: 15px;
                padding: 10px;
            }}
            .slot {{
                width: {slot_size};
                height: 55px;
                background: linear-gradient(145deg, #fff, #f5f5f5);
                border-radius: 10px;
                display: flex;
                align-items: center;
                justify-content: center;
                font-size: {font_size};
                font-weight: 800;
                color: #333;
                box-shadow: 0 4px 15px rgba(0,0,0,0.1);
                border: 2px solid #FF9800;
                overflow: hidden;
                position: relative;
                opacity: 0;
                transform: scale(0.8);
                transition: all 0.3s ease;
            }}
            .slot.active {{
                opacity: 1;
                transform: scale(1);
                animation: glow 0.08s infinite alternate;
            }}
            .slot.winner {{
                background: linear-gradient(135deg, #FF9800, #FF5722);
                color: white;
                border-color: #E65100;
                animation: popIn 0.4s ease;
                opacity: 1;
                transform: scale(1);
            }}
            @keyframes glow {{
                0% {{ box-shadow: 0 0 5px #FF9800, inset 0 0 5px rgba(255,152,0,0.2); }}
                100% {{ box-shadow: 0 0 15px #FF9800, inset 0 0 10px rgba(255,152,0,0.3); }}
            }}
            @keyframes popIn {{
                0% {{ transform: scale(0.5); }}
                50% {{ transform: scale(1.15); }}
                100% {{ transform: scale(1); }}
            }}
            .progress-container {{
                width: 100%;
                margin: 10px 0;
            }}
            .progress {{
                width: 100%;
                height: 8px;
                background: #e0e0e0;
                border-radius: 5px;
                overflow: hidden;
            }}
            .progress-bar {{
                height: 100%;
                background: linear-gradient(90deg, #4CAF50, #8BC34A);
                width: 0%;
                transition: width 0.15s ease;
                border-radius: 5px;
            }}
            .status {{
                font-size: 1.3rem;
                font-weight: 700;
                color: #4CAF50;
                margin-top: 10px;
                padding: 12px 25px;
                background: linear-gradient(145deg, #fff, #f8f9fa);
                border-radius: 10px;
                border: 2px solid #4CAF50;
                display: none;
                animation: popIn 0.5s ease;
            }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <div class="prize-title">🎲 {prize_name}</div>
                <div class="counter" id="counter">Mengundi 0/{total_winners} pemenang...</div>
            </div>
            <div class="slots-container" id="slotsContainer"></div>
            <div class="progress-container">
                <div class="progress"><div class="progress-bar" id="progressBar"></div></div>
            </div>
            <div class="status" id="status">🎉 {total_winners} PEMENANG TERPILIH!</div>
        </div>
        <script>
            const winners = {winners_js};
            const allNums = {all_nums_js};
            const container = document.getElementById('slotsContainer');
            const progressBar = document.getElementById('progressBar');
            const status = document.getElementById('status');
            const counter = document.getElementById('counter');
            const totalWinners = winners.length;
            
            // Create all slot elements (hidden initially)
            winners.forEach((_, idx) => {{
                const slot = document.createElement('div');
                slot.className = 'slot';
                slot.id = 'slot' + idx;
                slot.innerHTML = '<span class="slot-number">????</span>';
                container.appendChild(slot);
            }});
            
            function getRandomNum() {{
                return allNums[Math.floor(Math.random() * allNums.length)];
            }}
            
            let revealedCount = 0;
            const baseDelay = 80; // ms between each winner reveal
            const spinDuration = 600; // ms for spin animation per slot
            
            function revealWinner(slotIdx) {{
                const slot = document.getElementById('slot' + slotIdx);
                const numSpan = slot.querySelector('.slot-number');
                
                // Make slot visible and start spinning
                slot.classList.add('active');
                
                let spinCount = 0;
                const maxSpins = Math.floor(spinDuration / 40);
                
                function spin() {{
                    if (spinCount < maxSpins) {{
                        numSpan.textContent = getRandomNum();
                        spinCount++;
                        setTimeout(spin, 40);
                    }} else {{
                        // Reveal winner
                        slot.classList.remove('active');
                        slot.classList.add('winner');
                        numSpan.textContent = winners[slotIdx];
                        revealedCount++;
                        
                        // Update counter and progress
                        counter.textContent = 'Mengundi ' + revealedCount + '/{total_winners} pemenang...';
                        progressBar.style.width = (revealedCount / totalWinners * 100) + '%';
                        
                        if (revealedCount === totalWinners) {{
                            counter.textContent = '✅ Selesai!';
                            status.style.display = 'block';
                        }}
                    }}
                }}
                spin();
            }}
            
            // Cascade reveal - start each winner after a delay
            setTimeout(() => {{
                winners.forEach((_, idx) => {{
                    setTimeout(() => revealWinner(idx), idx * baseDelay);
                }});
            }}, 300);
        </script>
    </body>
    </html>
    '''
    return html

def create_spinning_wheel_html(all_participants, winner, wheel_size=400):
    """Create a clean, focused lottery animation"""
    total_pool = len(all_participants)
    all_nums_js = json.dumps(all_participants)
    
    html = f'''
    <!DOCTYPE html>
    <html>
    <head>
        <style>
            * {{ margin: 0; padding: 0; box-sizing: border-box; }}
            body {{
                display: flex;
                align-items: center;
                justify-content: center;
                min-height: 100vh;
                background: transparent;
                font-family: 'Segoe UI', sans-serif;
            }}
            .container {{
                text-align: center;
                width: 100%;
                max-width: 400px;
            }}
            .pool-badge {{
                background: linear-gradient(135deg, #667eea, #764ba2);
                padding: 8px 20px;
                border-radius: 20px;
                display: inline-block;
                margin-bottom: 15px;
            }}
            .pool-badge p {{
                color: white;
                margin: 0;
                font-size: 0.9rem;
            }}
            .pool-badge .total {{
                font-weight: bold;
            }}
            .number-display {{
                background: linear-gradient(145deg, #1a1a2e, #16213e);
                border-radius: 20px;
                padding: 40px 30px;
                box-shadow: 0 15px 40px rgba(0,0,0,0.4);
            }}
            .spinning-number {{
                font-size: 5rem;
                font-weight: 900;
                color: #00ff88;
                text-shadow: 0 0 30px rgba(0,255,136,0.6);
                font-family: 'Courier New', monospace;
                letter-spacing: 8px;
            }}
            .status-text {{
                color: #aaa;
                font-size: 1rem;
                margin-top: 15px;
            }}
            .progress {{
                width: 80%;
                height: 8px;
                background: #333;
                border-radius: 4px;
                margin: 20px auto 0;
                overflow: hidden;
            }}
            .progress-bar {{
                height: 100%;
                background: linear-gradient(90deg, #00ff88, #00ccff, #E91E63);
                width: 0%;
                transition: width 0.1s;
            }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="pool-badge">
                <p>Mengundi dari <span class="total">{total_pool}</span> peserta</p>
            </div>
            
            <div class="number-display">
                <div class="spinning-number" id="spinNumber">----</div>
                <div class="status-text" id="statusText">Mengacak...</div>
                <div class="progress"><div class="progress-bar" id="progressBar"></div></div>
            </div>
        </div>
        
        <script>
            const allNums = {all_nums_js};
            const winner = "{winner}";
            const totalPool = {total_pool};
            const spinNumber = document.getElementById('spinNumber');
            const statusText = document.getElementById('statusText');
            const progressBar = document.getElementById('progressBar');
            
            const duration = 4000;
            const startTime = Date.now();
            
            function getRandomNum() {{
                return allNums[Math.floor(Math.random() * totalPool)];
            }}
            
            function animate() {{
                const elapsed = Date.now() - startTime;
                const progress = Math.min(elapsed / duration, 1);
                progressBar.style.width = (progress * 100) + '%';
                
                if (progress < 0.85) {{
                    spinNumber.textContent = getRandomNum();
                    const delay = progress < 0.6 ? 40 : 60 + (progress - 0.6) * 400;
                    setTimeout(animate, delay);
                }} else if (progress < 1) {{
                    spinNumber.textContent = getRandomNum();
                    statusText.textContent = 'Hampir selesai...';
                    setTimeout(animate, 120);
                }} else {{
                    spinNumber.textContent = winner;
                    spinNumber.style.color = '#FFD700';
                    spinNumber.style.textShadow = '0 0 40px rgba(255,215,0,0.8)';
                    statusText.textContent = '✅ Selesai!';
                }}
            }}
            
            setTimeout(animate, 200);
        </script>
    </body>
    </html>
    '''
    return html

def generate_pptx(results_df, prize_tiers):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    for tier in prize_tiers:
        tier_winners = results_df[results_df["Hadiah"] == tier["name"]].copy()
        if len(tier_winners) == 0:
            continue
        
        tier_winners = tier_winners.sort_values(by="Nomor Undian", ascending=True).reset_index(drop=True)
        
        cols = 5
        rows_per_slide = 5
        winners_per_slide = cols * rows_per_slide
        total_winners = len(tier_winners)
        num_slides = (total_winners + winners_per_slide - 1) // winners_per_slide
        
        for slide_num in range(num_slides):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            background.fill.gradient()
            background.fill.gradient_stops[0].color.rgb = RGBColor(245, 87, 108)
            background.fill.gradient_stops[1].color.rgb = RGBColor(240, 147, 251)
            background.line.fill.background()
            
            title_text = f"{tier['icon']} {tier['name']}"
            if num_slides > 1:
                title_text += f" ({slide_num + 1}/{num_slides})"
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title_text
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            
            cell_width = Inches(2.4)
            cell_height = Inches(1.1)
            gap_x = Inches(0.1)
            gap_y = Inches(0.1)
            
            total_grid_width = cols * cell_width + (cols - 1) * gap_x
            start_x = (prs.slide_width - total_grid_width) / 2
            start_y = Inches(1.3)
            
            start_idx = slide_num * winners_per_slide
            end_idx = min(start_idx + winners_per_slide, total_winners)
            
            for idx, (_, row) in enumerate(tier_winners.iloc[start_idx:end_idx].iterrows()):
                row_num = idx // cols
                col_num = idx % cols
                
                left = start_x + col_num * (cell_width + gap_x)
                top = start_y + row_num * (cell_height + gap_y)
                
                shape = slide.shapes.add_shape(5, left, top, cell_width, cell_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(245, 87, 108)
                shape.line.width = Pt(2)
                
                nomor = str(row["Nomor Undian"])
                nama_raw = row.get("Nama", "")
                nama = str(nama_raw) if pd.notna(nama_raw) else "-"
                if nama.lower() == "nan":
                    nama = "-"
                hp_raw = row.get("No HP", "")
                hp = format_phone(hp_raw)
                
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.space_before = Pt(4)
                p.space_after = Pt(0)
                run = p.add_run()
                run.text = nomor
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(51, 51, 51)
                
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(2)
                p2.space_after = Pt(0)
                run2 = p2.add_run()
                run2.text = nama
                run2.font.size = Pt(12)
                run2.font.color.rgb = RGBColor(102, 102, 102)
                
                p3 = tf.add_paragraph()
                p3.alignment = PP_ALIGN.CENTER
                p3.space_before = Pt(0)
                run3 = p3.add_run()
                run3.text = hp
                run3.font.size = Pt(11)
                run3.font.color.rgb = RGBColor(136, 136, 136)
    
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()

def generate_shuffle_pptx(winners_list, prize_name, name_lookup=None, phone_lookup=None):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    if name_lookup is None:
        name_lookup = {}
    if phone_lookup is None:
        phone_lookup = {}
    
    sorted_winners = sorted(winners_list, key=lambda x: str(x))
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.gradient()
    background.fill.gradient_stops[0].color.rgb = RGBColor(255, 152, 0)
    background.fill.gradient_stops[1].color.rgb = RGBColor(255, 87, 34)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"🎲 {prize_name}"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    cols = 5
    cell_width = Inches(2.4)
    cell_height = Inches(1.1)
    gap_x = Inches(0.1)
    gap_y = Inches(0.1)
    
    total_grid_width = cols * cell_width + (cols - 1) * gap_x
    start_x = (prs.slide_width - total_grid_width) / 2
    start_y = Inches(1.3)
    
    for idx, winner in enumerate(sorted_winners):
        row_num = idx // cols
        col_num = idx % cols
        
        left = start_x + col_num * (cell_width + gap_x)
        top = start_y + row_num * (cell_height + gap_y)
        
        shape = slide.shapes.add_shape(5, left, top, cell_width, cell_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(255, 152, 0)
        shape.line.width = Pt(2)
        
        nomor = str(winner)
        nama_raw = name_lookup.get(winner, "")
        nama = str(nama_raw) if pd.notna(nama_raw) else "-"
        if nama.lower() == "nan":
            nama = "-"
        hp = format_phone(phone_lookup.get(winner, ""))
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = nomor
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(51, 51, 51)
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        p2.space_after = Pt(0)
        run2 = p2.add_run()
        run2.text = nama
        run2.font.size = Pt(12)
        run2.font.color.rgb = RGBColor(102, 102, 102)
        
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(0)
        run3 = p3.add_run()
        run3.text = hp
        run3.font.size = Pt(11)
        run3.font.color.rgb = RGBColor(136, 136, 136)
    
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()

def generate_shuffle_pptx_v2(prize_assignments, name_lookup=None, phone_lookup=None, session_name="Sesi"):
    """Generate PPT with one slide per prize category"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    if name_lookup is None:
        name_lookup = {}
    if phone_lookup is None:
        phone_lookup = {}
    
    # Group by prize
    prize_groups = {}
    for pa in prize_assignments:
        prize = pa["prize"]
        if prize not in prize_groups:
            prize_groups[prize] = []
        prize_groups[prize].append(pa["winner"])
    
    slide_layout = prs.slide_layouts[6]
    
    for prize_name, winners in prize_groups.items():
        sorted_winners = sorted(winners, key=lambda x: str(x))
        
        # 15 winners per slide (5x3 grid)
        winners_per_slide = 15
        total_slides = (len(sorted_winners) + winners_per_slide - 1) // winners_per_slide
        
        for slide_num in range(total_slides):
            slide = prs.slides.add_slide(slide_layout)
            
            # Background
            background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            background.fill.gradient()
            background.fill.gradient_stops[0].color.rgb = RGBColor(76, 175, 80)
            background.fill.gradient_stops[1].color.rgb = RGBColor(56, 142, 60)
            background.line.fill.background()
            
            # Title with prize name
            title_text = f"🎁 {prize_name}"
            if total_slides > 1:
                title_text += f" ({slide_num + 1}/{total_slides})"
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.6))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title_text
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            
            # Subtitle
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.33), Inches(0.4))
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = f"{session_name} - {len(sorted_winners)} Pemenang"
            run2.font.size = Pt(18)
            run2.font.color.rgb = RGBColor(255, 255, 255)
            
            # Grid layout (5x3 = 15 per slide)
            cols = 5
            rows_per_slide = 3
            cell_width = Inches(2.4)
            cell_height = Inches(1.6)
            gap_x = Inches(0.1)
            gap_y = Inches(0.15)
            
            total_grid_width = cols * cell_width + (cols - 1) * gap_x
            start_x = (prs.slide_width - total_grid_width) / 2
            start_y = Inches(1.3)
            
            start_idx = slide_num * winners_per_slide
            end_idx = min(start_idx + winners_per_slide, len(sorted_winners))
            slide_winners = sorted_winners[start_idx:end_idx]
            
            for idx, winner in enumerate(slide_winners):
                row_num = idx // cols
                col_num = idx % cols
                
                left = start_x + col_num * (cell_width + gap_x)
                top = start_y + row_num * (cell_height + gap_y)
                
                shape = slide.shapes.add_shape(5, left, top, cell_width, cell_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(76, 175, 80)
                shape.line.width = Pt(2)
                
                nomor = str(winner)
                nama_raw = name_lookup.get(winner, "")
                nama = str(nama_raw) if pd.notna(nama_raw) else "-"
                if nama.lower() == "nan":
                    nama = "-"
                hp = format_phone(phone_lookup.get(winner, ""))
                
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.space_before = Pt(4)
                p.space_after = Pt(0)
                run = p.add_run()
                run.text = nomor
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(51, 51, 51)
                
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(2)
                p2.space_after = Pt(0)
                run2 = p2.add_run()
                run2.text = nama
                run2.font.size = Pt(12)
                run2.font.color.rgb = RGBColor(102, 102, 102)
                
                p3 = tf.add_paragraph()
                p3.alignment = PP_ALIGN.CENTER
                p3.space_before = Pt(0)
                run3 = p3.add_run()
                run3.text = hp
                run3.font.size = Pt(11)
                run3.font.color.rgb = RGBColor(136, 136, 136)
    
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()

def generate_wheel_pptx(winners_list, prizes_list, name_lookup=None, phone_lookup=None):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    if name_lookup is None:
        name_lookup = {}
    if phone_lookup is None:
        phone_lookup = {}
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.gradient()
    background.fill.gradient_stops[0].color.rgb = RGBColor(233, 30, 99)
    background.fill.gradient_stops[1].color.rgb = RGBColor(156, 39, 176)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "🎡 PEMENANG SPINNING WHEEL"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    start_y = Inches(1.5)
    
    for idx, (winner, prize) in enumerate(zip(winners_list, prizes_list)):
        nama_raw = name_lookup.get(winner, "")
        nama = str(nama_raw) if pd.notna(nama_raw) else "-"
        if nama.lower() == "nan":
            nama = "-"
        hp = format_phone(phone_lookup.get(winner, ""))
        
        text_box = slide.shapes.add_textbox(Inches(0.5), start_y + Inches(idx * 0.55), Inches(12.33), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"#{idx+1}: {winner} - {nama} ({hp}) | {prize}"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()

st.set_page_config(page_title="Undian Move & Groove", layout="wide", initial_sidebar_state="collapsed")

st.markdown("""
<style>
    .main-title { text-align: center; color: white; font-size: 2.5rem; font-weight: 800; margin: 0; }
    .subtitle { text-align: center; color: #ccc; font-size: 1.2rem; margin: 0; }
    .section-header { text-align: center; color: white; font-size: 1.5rem; font-weight: 700; margin: 1rem 0; }
    .stats-card { background: rgba(255,255,255,0.1); border-radius: 15px; padding: 1.5rem; text-align: center; }
    .stats-number { font-size: 2.5rem; font-weight: 800; color: #f5576c; }
    .stats-label { color: #ccc; font-size: 1rem; }
    .prize-header { background: linear-gradient(135deg, #f5576c, #f093fb); padding: 2rem; border-radius: 20px; text-align: center; color: white; margin: 1rem 0; }
    .winner-btn { background: linear-gradient(135deg, #4CAF50, #8BC34A); color: white; padding: 1rem; border-radius: 10px; text-align: center; margin: 0.5rem 0; cursor: pointer; }
</style>
""", unsafe_allow_html=True)

if "current_page" not in st.session_state:
    st.session_state["current_page"] = "home"
if "prize_tiers" not in st.session_state:
    saved_config = load_prize_config()
    st.session_state["prize_tiers"] = saved_config if saved_config else PRIZE_TIERS.copy()

# Auto-load saved lottery results on startup
if "results_loaded" not in st.session_state:
    if load_lottery_results():
        st.session_state["results_loaded"] = True
        st.toast("✅ Hasil undian sebelumnya berhasil dimuat!", icon="💾")
    else:
        st.session_state["results_loaded"] = True

if os.path.exists("attached_assets/Small Banner-01_1764081768006.png"):
    st.image("attached_assets/Small Banner-01_1764081768006.png", use_container_width=True)
st.markdown('<p class="main-title">🎉 UNDIAN MOVE & GROOVE 🎉</p>', unsafe_allow_html=True)
st.markdown('<p class="subtitle">7 Desember 2024</p>', unsafe_allow_html=True)

# Show backup status indicator and reset button
evoucher_done = st.session_state.get("evoucher_done", False)
shuffle_done = st.session_state.get("shuffle_done", False)
wheel_done = st.session_state.get("wheel_done", False)
current_file = st.session_state.get("current_results_file", "")

if evoucher_done or shuffle_done or wheel_done or current_file:
    status_parts = []
    if evoucher_done:
        status_parts.append("E-Voucher ✓")
    if st.session_state.get("shuffle_results", {}):
        status_parts.append(f"Shuffle ({len(st.session_state.get('shuffle_results', {}))}/3) ✓")
    if st.session_state.get("wheel_winners", []):
        status_parts.append(f"Wheel ({len(st.session_state.get('wheel_winners', []))}/10) ✓")
    
    col_status, col_reset = st.columns([4, 1])
    
    with col_status:
        if status_parts:
            status_text = " | ".join(status_parts)
            file_info = f"📁 {current_file}" if current_file else ""
            gdrive_status = st.session_state.get("gdrive_save_status", None)
            gdrive_icon = "☁️✓" if gdrive_status else "💾"
            st.markdown(f"""
            <div style="background: rgba(76, 175, 80, 0.2); border: 1px solid #4CAF50; border-radius: 8px; padding: 0.5rem; text-align: center;">
                <span style="color: #4CAF50; font-size: 0.9rem;">{gdrive_icon} Auto-Save: {status_text}</span>
                <br><span style="color: #888; font-size: 0.75rem;">{file_info} {'| Google Drive ✓' if gdrive_status else ''}</span>
            </div>
            """, unsafe_allow_html=True)
    
    with col_reset:
        if st.button("🔄 RESET", key="reset_lottery", use_container_width=True, type="secondary"):
            reset_lottery_session()
            st.rerun()

current_page = st.session_state.get("current_page", "home")

if current_page == "home":
    tab1, tab2 = st.tabs(["📁 Upload File CSV", "🔗 Google Sheets URL"])
    
    df = None
    
    with tab1:
        uploaded_file = st.file_uploader("Upload File CSV", type=["csv"], help="File CSV harus berisi kolom 'Nomor Undian'")
        if uploaded_file:
            try:
                uploaded_file.seek(0)
                file_content = uploaded_file.read()
                uploaded_file.seek(0)
                
                content_hash = hashlib.md5(file_content).hexdigest()
                if st.session_state.get("last_content_hash") != content_hash:
                    st.session_state["last_content_hash"] = content_hash
                    st.session_state["data_source_changed"] = True
                    st.session_state["evoucher_done"] = False
                    st.session_state["evoucher_results"] = None
                    st.session_state["shuffle_results"] = {}
                    st.session_state["shuffle_done"] = False
                    st.session_state["wheel_winners"] = []
                    st.session_state["wheel_prizes"] = []
                    st.session_state["wheel_done"] = False
                    if "sheets_df" in st.session_state:
                        del st.session_state["sheets_df"]
                    if "last_sheets_hash" in st.session_state:
                        del st.session_state["last_sheets_hash"]
                    if "remaining_pool" in st.session_state:
                        del st.session_state["remaining_pool"]
                
                df = pd.read_csv(uploaded_file, dtype=str, encoding='utf-8-sig')
                df.columns = df.columns.str.strip().str.replace('\ufeff', '')
            except Exception as e:
                st.error(f"Error: {e}")
    
    with tab2:
        st.markdown("""
        <div style="background: rgba(76, 175, 80, 0.2); border: 1px solid #4CAF50; border-radius: 8px; padding: 0.8rem; margin-bottom: 1rem;">
            <p style="color: #4CAF50; margin: 0; font-size: 0.9rem;">
                ✅ <strong>Default:</strong> Google Sheets Move & Groove 7 Desember sudah terkonfigurasi
            </p>
        </div>
        """, unsafe_allow_html=True)
        
        sheets_url = st.text_input("Google Sheets URL", value=DEFAULT_SHEETS_URL, help="URL sudah diisi otomatis dengan data Move & Groove")
        
        col_load, col_refresh = st.columns(2)
        with col_load:
            load_btn = st.button("📥 Ambil Data", use_container_width=True)
        with col_refresh:
            refresh_btn = st.button("🔄 Refresh Data", use_container_width=True)
        
        if (load_btn or refresh_btn) and sheets_url:
            try:
                sheet_id_match = re.search(r'/d/([a-zA-Z0-9-_]+)', sheets_url)
                if sheet_id_match:
                    sheet_id = sheet_id_match.group(1)
                    
                    gid_match = re.search(r'gid=(\d+)', sheets_url)
                    gid = gid_match.group(1) if gid_match else "0"
                    
                    csv_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}"
                    response = requests.get(csv_url, timeout=30)
                    response.raise_for_status()
                    
                    content_hash = hashlib.md5(response.content).hexdigest()
                    
                    if refresh_btn or st.session_state.get("last_sheets_hash") != content_hash:
                        if not refresh_btn:
                            st.session_state["last_sheets_hash"] = content_hash
                            st.session_state["data_source_changed"] = True
                            st.session_state["evoucher_done"] = False
                            st.session_state["evoucher_results"] = None
                            st.session_state["shuffle_results"] = {}
                            st.session_state["shuffle_done"] = False
                            st.session_state["wheel_winners"] = []
                            st.session_state["wheel_prizes"] = []
                            st.session_state["wheel_done"] = False
                            if "last_content_hash" in st.session_state:
                                del st.session_state["last_content_hash"]
                            if "remaining_pool" in st.session_state:
                                del st.session_state["remaining_pool"]
                    
                    df = pd.read_csv(StringIO(response.content.decode('utf-8-sig')), dtype=str)
                    df.columns = df.columns.str.strip().str.replace('\ufeff', '')
                    st.session_state["sheets_df"] = df
                    st.session_state["last_sheets_hash"] = content_hash
                    st.success(f"✅ Berhasil mengambil {len(df)} baris data dari Google Sheets!")
            except Exception as e:
                st.error(f"Error: {e}")
        
        if df is None and "sheets_df" in st.session_state:
            df = st.session_state["sheets_df"]
    
    if df is not None:
        undian_col = None
        for col in df.columns:
            if "undian" in col.lower():
                undian_col = col
                break
        
        if undian_col is None:
            st.error("❌ File harus memiliki kolom 'Nomor Undian'")
        else:
            df = df.rename(columns={undian_col: "Nomor Undian"})
            
            name_col = None
            for col in df.columns:
                if "nama" in col.lower():
                    name_col = col
                    break
            if name_col and name_col != "Nama":
                df = df.rename(columns={name_col: "Nama"})
            elif "Nama" not in df.columns:
                df["Nama"] = ""
            
            phone_col = None
            for col in df.columns:
                if "hp" in col.lower() or "phone" in col.lower() or "telepon" in col.lower():
                    phone_col = col
                    break
            if phone_col and phone_col != "No HP":
                df = df.rename(columns={phone_col: "No HP"})
            elif "No HP" not in df.columns:
                df["No HP"] = ""
            
            df["Nomor Undian"] = df["Nomor Undian"].astype(str).str.strip().str.zfill(4)
            df = df.dropna(subset=["Nomor Undian"])
            df = df[df["Nomor Undian"].str.len() > 0]
            
            df["Eligible"] = df.apply(lambda x: is_eligible_for_prize(x.get("Nama", ""), x.get("No HP", "")), axis=1)
            
            st.session_state["participant_data"] = df
            eligible_df = df[df["Eligible"] == True]
            st.session_state["eligible_participants"] = eligible_df["Nomor Undian"].tolist()
            
            if "remaining_pool" not in st.session_state or st.session_state.get("data_source_changed", False):
                st.session_state["remaining_pool"] = eligible_df.copy()
                st.session_state["data_source_changed"] = False
            
            total_all = len(df)
            total_eligible = len(eligible_df)
            total_excluded = total_all - total_eligible
            remaining_pool = st.session_state.get("remaining_pool", eligible_df)
            
            st.success(f"✅ Data: {total_all} peserta ({total_eligible} eligible, {total_excluded} VIP/F) | Sisa: {len(remaining_pool)}")
            
            st.markdown("<br>", unsafe_allow_html=True)
            st.markdown("<p style='text-align:center; color:white; font-size:1.8rem; font-weight:bold;'>🎯 PILIH JENIS UNDIAN</p>", unsafe_allow_html=True)
            st.markdown("<br>", unsafe_allow_html=True)
            
            evoucher_done = st.session_state.get("evoucher_done", False)
            shuffle_done = st.session_state.get("shuffle_done", False)
            wheel_done = st.session_state.get("wheel_done", False)
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                prize_tiers = st.session_state.get("prize_tiers", PRIZE_TIERS)
                total_prizes = calculate_total_winners(prize_tiers)
                status_text = "✅ SELESAI" if evoucher_done else f"{total_prizes} hadiah"
                status_color = "#4CAF50" if evoucher_done else "white"
                st.markdown(f"""
                <div style="background: rgba(76,175,80,0.2); border-radius: 15px; padding: 1.5rem; text-align: center; border: 2px solid #4CAF50; min-height: 200px;">
                    <p style="color: #4CAF50; font-size: 1.8rem; font-weight: bold; margin: 0;">🎁 E-Voucher</p>
                    <p style="color: {status_color}; font-size: 1.1rem; margin: 0.5rem 0;">{status_text}</p>
                    <p style="color: #aaa; font-size: 0.9rem; margin: 0.5rem 0;">
                    4 Kategori Voucher<br>
                    Tokopedia, Indomaret, Bensin, SNL
                    </p>
                </div>
                """, unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("🎁 UNDIAN E-VOUCHER", key="btn_evoucher", use_container_width=True):
                    st.session_state["current_page"] = "evoucher_page"
                    st.rerun()
            
            with col2:
                shuffle_results = st.session_state.get("shuffle_results", {})
                completed_sessions = len(shuffle_results)
                status_text = f"✅ {completed_sessions}/3 Sesi" if completed_sessions > 0 else "3 Sesi x 30 hadiah"
                st.markdown(f"""
                <div style="background: rgba(255,152,0,0.2); border-radius: 15px; padding: 1.5rem; text-align: center; border: 2px solid #FF9800; min-height: 200px;">
                    <p style="color: #FF9800; font-size: 1.8rem; font-weight: bold; margin: 0;">🎲 Shuffle</p>
                    <p style="color: white; font-size: 1.1rem; margin: 0.5rem 0;">{status_text}</p>
                    <p style="color: #aaa; font-size: 0.9rem; margin: 0.5rem 0;">
                    Lucky Draw 3 Sesi<br>
                    Sisa: {len(remaining_pool)} peserta
                    </p>
                </div>
                """, unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("🎲 UNDIAN SHUFFLE", key="btn_shuffle", use_container_width=True):
                    st.session_state["current_page"] = "shuffle_page"
                    st.rerun()
            
            with col3:
                wheel_winners = st.session_state.get("wheel_winners", [])
                status_text = f"✅ {len(wheel_winners)}/10 Hadiah" if len(wheel_winners) > 0 else "10 Hadiah Utama"
                st.markdown(f"""
                <div style="background: rgba(233,30,99,0.2); border-radius: 15px; padding: 1.5rem; text-align: center; border: 2px solid #E91E63; min-height: 200px;">
                    <p style="color: #E91E63; font-size: 1.8rem; font-weight: bold; margin: 0;">🎡 Spinning Wheel</p>
                    <p style="color: white; font-size: 1.1rem; margin: 0.5rem 0;">{status_text}</p>
                    <p style="color: #aaa; font-size: 0.9rem; margin: 0.5rem 0;">
                    Grand Prize satu per satu<br>
                    Sisa: {len(remaining_pool)} peserta
                    </p>
                </div>
                """, unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("🎡 UNDIAN WHEEL", key="btn_wheel", use_container_width=True):
                    st.session_state["current_page"] = "wheel_page"
                    st.rerun()
            
            evoucher_results = st.session_state.get("evoucher_results")
            shuffle_results = st.session_state.get("shuffle_results", {})
            wheel_winners = st.session_state.get("wheel_winners", [])
            
            has_results = (evoucher_results is not None) or (len(shuffle_results) > 0) or (len(wheel_winners) > 0)
            
            if has_results:
                st.markdown("<br>", unsafe_allow_html=True)
                st.markdown("---")
                st.markdown("<p style='text-align:center; color:white; font-size:1.5rem; font-weight:bold;'>🏆 HASIL PEMENANG</p>", unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)
                
                if evoucher_results is not None:
                    st.markdown("<p style='color:#4CAF50; font-size:1.2rem; font-weight:bold;'>🎁 E-Voucher (4 Kategori)</p>", unsafe_allow_html=True)
                    prize_tiers = st.session_state.get("prize_tiers", PRIZE_TIERS)
                    cols = st.columns(4)
                    for idx, tier in enumerate(prize_tiers):
                        with cols[idx]:
                            tier_winners = evoucher_results[evoucher_results["Hadiah"] == tier["name"]]
                            count = len(tier_winners)
                            if st.button(f"{tier['icon']} {tier['name'].split()[0]}\n({count})", key=f"home_ev_{idx}", use_container_width=True):
                                st.session_state["viewing_tier"] = tier
                                st.session_state["current_page"] = "evoucher_category"
                                st.rerun()
                    st.markdown("<br>", unsafe_allow_html=True)
                
                if len(shuffle_results) > 0:
                    st.markdown("<p style='color:#FF9800; font-size:1.2rem; font-weight:bold;'>🎲 Shuffle (3 Sesi)</p>", unsafe_allow_html=True)
                    cols = st.columns(3)
                    for i in range(3):
                        with cols[i]:
                            batch_key = f"shuffle_batch_{i}"
                            if batch_key in shuffle_results:
                                result = shuffle_results[batch_key]
                                prize_name = result.get("prize_name", f"Sesi {i+1}")
                                count = len(result.get("winners", []))
                                if st.button(f"🎲 Sesi {i+1}\n({count})", key=f"home_sh_{i}", use_container_width=True):
                                    st.session_state["viewing_shuffle_batch"] = i
                                    st.session_state["current_page"] = "shuffle_results"
                                    st.rerun()
                            else:
                                st.button(f"🎲 Sesi {i+1}\n(Belum)", key=f"home_sh_{i}", use_container_width=True, disabled=True)
                    st.markdown("<br>", unsafe_allow_html=True)
                
                if len(wheel_winners) > 0:
                    st.markdown("<p style='color:#E91E63; font-size:1.2rem; font-weight:bold;'>🎡 Spinning Wheel</p>", unsafe_allow_html=True)
                    if st.button(f"🎡 Grand Prize ({len(wheel_winners)} pemenang)", key="home_wheel", use_container_width=True):
                        st.session_state["current_page"] = "wheel_results"
                        st.rerun()
            
            st.markdown("<br>", unsafe_allow_html=True)
            st.markdown("---")
            if st.button("🔄 RESET UNDIAN (Mulai dari Awal)", key="reset_all", use_container_width=True):
                keys_to_keep = ["prize_tiers", "participant_data", "eligible_participants"]
                for key in list(st.session_state.keys()):
                    if key not in keys_to_keep:
                        del st.session_state[key]
                st.session_state["remaining_pool"] = eligible_df.copy()
                st.session_state["current_page"] = "home"
                st.rerun()
    
    else:
        st.markdown("<br>", unsafe_allow_html=True)
        st.info("📁 Silakan upload file CSV atau paste URL Google Sheets untuk memulai undian.")

elif current_page == "evoucher_page":
    prize_tiers = st.session_state.get("prize_tiers", PRIZE_TIERS)
    total_prizes = calculate_total_winners(prize_tiers)
    evoucher_results = st.session_state.get("evoucher_results")
    
    if st.button("⬅️ KEMBALI KE MENU", key="back_to_home"):
        st.session_state["current_page"] = "home"
        st.rerun()
    
    st.markdown(f"""
    <div style="background: linear-gradient(135deg, #4CAF50, #8BC34A); padding: 2rem; border-radius: 15px; text-align: center; margin: 1rem 0;">
        <p style="color: white; font-size: 2.5rem; font-weight: bold; margin: 0;">🎁 UNDIAN E-VOUCHER</p>
        <p style="color: #fff; font-size: 1.2rem; margin: 0.5rem 0;">{total_prizes} Hadiah - {len(prize_tiers)} Kategori</p>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("<p style='text-align:center; color:white; font-size:1.3rem; font-weight:bold;'>⚙️ KONFIGURASI HADIAH</p>", unsafe_allow_html=True)
    
    with st.expander("📝 Edit Jenis & Jumlah Hadiah", expanded=evoucher_results is None):
        new_tiers = []
        cols = st.columns(4)
        for idx, tier in enumerate(prize_tiers):
            with cols[idx % 4]:
                st.markdown(f"**{tier['icon']} Kategori {idx+1}**")
                name = st.text_input(f"Nama", value=tier["name"], key=f"tier_name_{idx}")
                count = st.number_input(f"Jumlah", value=tier["count"], min_value=1, max_value=500, key=f"tier_count_{idx}")
                icon = st.text_input(f"Icon", value=tier["icon"], key=f"tier_icon_{idx}")
                new_tiers.append({"name": name, "icon": icon, "count": count})
        
        if st.button("💾 Simpan Konfigurasi", use_container_width=True):
            start = 1
            for tier in new_tiers:
                tier["start"] = start
                tier["end"] = start + tier["count"] - 1
                start = tier["end"] + 1
            st.session_state["prize_tiers"] = new_tiers
            save_prize_config(new_tiers)
            st.success("✅ Konfigurasi disimpan!")
            st.rerun()
    
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("<p style='text-align:center; color:white; font-size:1.3rem; font-weight:bold;'>📋 KATEGORI HADIAH</p>", unsafe_allow_html=True)
    
    cols = st.columns(4)
    for idx, tier in enumerate(prize_tiers):
        with cols[idx % 4]:
            st.markdown(f"""
            <div style="background: white; border-radius: 20px; padding: 1.5rem; text-align: center; margin-bottom: 1rem; box-shadow: 0 4px 15px rgba(0,0,0,0.2);">
                <div style="font-size: 3rem; margin-bottom: 0.5rem;">{tier['icon']}</div>
                <p style="color: #333; font-size: 1rem; font-weight: bold; margin: 0;">{tier['name']}</p>
                <p style="color: #f5576c; font-size: 1.2rem; font-weight: bold; margin: 0.3rem 0;">{tier['count']} Pemenang</p>
            </div>
            """, unsafe_allow_html=True)
    
    if evoucher_results is None:
        st.markdown("<br>", unsafe_allow_html=True)
        
        if st.button("🎲 MULAI UNDIAN E-VOUCHER", key="start_evoucher", use_container_width=True):
            eligible_participants = st.session_state.get("eligible_participants", [])
            
            if len(eligible_participants) < total_prizes:
                st.error(f"❌ Peserta eligible ({len(eligible_participants)}) kurang dari total hadiah ({total_prizes})")
            else:
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                for i in range(100):
                    progress_bar.progress(i + 1)
                    if i < 30:
                        status_text.markdown(f"<p style='text-align:center; font-size:1.5rem; color:white;'>🔄 Mengumpulkan data... {i+1}%</p>", unsafe_allow_html=True)
                    elif i < 70:
                        status_text.markdown(f"<p style='text-align:center; font-size:1.5rem; color:white;'>🎲 Mengacak peserta... {i+1}%</p>", unsafe_allow_html=True)
                    else:
                        status_text.markdown(f"<p style='text-align:center; font-size:1.5rem; color:white;'>🏆 Menentukan pemenang... {i+1}%</p>", unsafe_allow_html=True)
                    time.sleep(0.02)
                
                shuffled = secure_shuffle(eligible_participants)
                winners = shuffled[:total_prizes]
                
                participant_data = st.session_state.get("participant_data")
                name_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["Nama"])) if participant_data is not None else {}
                phone_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["No HP"])) if participant_data is not None else {}
                
                results = []
                for i, winner in enumerate(winners, 1):
                    results.append({
                        "Peringkat": i,
                        "Nomor Undian": winner,
                        "Nama": name_lookup.get(winner, ""),
                        "No HP": phone_lookup.get(winner, ""),
                        "Hadiah": get_prize_dynamic(i, prize_tiers)
                    })
                
                results_df = pd.DataFrame(results)
                st.session_state["evoucher_results"] = results_df
                st.session_state["evoucher_done"] = True
                
                remaining = [p for p in eligible_participants if p not in winners]
                remaining_df = participant_data[participant_data["Nomor Undian"].isin(remaining)].copy() if participant_data is not None else pd.DataFrame()
                st.session_state["remaining_pool"] = remaining_df
                
                # Auto-save results
                save_lottery_results()
                
                progress_bar.empty()
                status_text.empty()
                st.balloons()
                st.rerun()
    
    else:
        st.markdown("<br>", unsafe_allow_html=True)
        st.success(f"✅ Undian E-Voucher selesai! {len(evoucher_results)} pemenang")
        
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown('<div class="section-header">🏆 LIHAT PEMENANG PER KATEGORI</div>', unsafe_allow_html=True)
        
        cols = st.columns(4)
        for idx, tier in enumerate(prize_tiers):
            with cols[idx % 4]:
                count = len(evoucher_results[evoucher_results["Hadiah"] == tier["name"]])
                if st.button(f"{tier['icon']} {tier['name']}\n({count} Pemenang)", key=f"view_tier_{idx}", use_container_width=True):
                    st.session_state["viewing_tier"] = tier
                    st.session_state["current_page"] = "evoucher_category"
                    st.rerun()
        
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("---")
        st.markdown("<p style='text-align:center; color:white; font-size:1.3rem; font-weight:600;'>📥 Download Hasil</p>", unsafe_allow_html=True)
        
        col1, col2 = st.columns(2)
        with col1:
            excel_buffer = BytesIO()
            with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                evoucher_results.to_excel(writer, index=False, sheet_name='Hasil Undian')
            st.download_button(
                label="📊 Download Excel (.xlsx)",
                data=excel_buffer.getvalue(),
                file_name="hasil_evoucher.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
        
        with col2:
            pptx_data = generate_pptx(evoucher_results, prize_tiers)
            st.download_button(
                label="📽️ Download PowerPoint (.pptx)",
                data=pptx_data,
                file_name="hasil_evoucher.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        
        st.markdown("<br>", unsafe_allow_html=True)
        remaining_pool = st.session_state.get("remaining_pool", pd.DataFrame())
        
        with st.expander(f"📋 Nomor yang Belum Diundi ({len(remaining_pool)} peserta)", expanded=False):
            if len(remaining_pool) > 0:
                remaining_numbers = remaining_pool["Nomor Undian"].tolist()
                cols = st.columns(15)
                for idx, num in enumerate(remaining_numbers[:150]):
                    with cols[idx % 15]:
                        st.markdown(f"<div style='background:#333;color:white;padding:0.3rem;border-radius:5px;text-align:center;margin:2px;font-size:0.8rem;'>{num}</div>", unsafe_allow_html=True)
                if len(remaining_numbers) > 150:
                    st.info(f"... dan {len(remaining_numbers) - 150} nomor lainnya")
            else:
                st.info("Semua nomor sudah diundi")
        
        if st.button("📊 SISA NOMOR → KEMBALI KE MENU UTAMA", key="ev_to_home", use_container_width=True):
            st.session_state["current_page"] = "home"
            st.rerun()

elif current_page == "evoucher_category":
    tier = st.session_state.get("viewing_tier")
    results_df = st.session_state.get("evoucher_results")
    
    if tier is None or results_df is None:
        st.session_state["current_page"] = "evoucher_page"
        st.rerun()
    
    if st.button("⬅️ KEMBALI", key="back_to_results"):
        st.session_state["current_page"] = "evoucher_page"
        st.rerun()
    
    tier_winners = results_df[results_df["Hadiah"] == tier["name"]].copy()
    tier_winners = tier_winners.sort_values(by="Nomor Undian", ascending=True).reset_index(drop=True)
    
    st.markdown(f"""
    <div class="prize-header">
        <div style="font-size: 4rem;">{tier["icon"]}</div>
        <div style="font-size: 2.5rem; font-weight: 800;">{tier["name"]}</div>
        <div style="font-size: 1.3rem;">{len(tier_winners)} Pemenang</div>
    </div>
    """, unsafe_allow_html=True)
    
    cols = 7
    rows = (len(tier_winners) + cols - 1) // cols
    
    for row in range(rows):
        row_cols = st.columns(cols)
        for col in range(cols):
            idx = row * cols + col
            if idx < len(tier_winners):
                winner = tier_winners.iloc[idx]
                with row_cols[col]:
                    nomor = winner["Nomor Undian"]
                    nama_raw = winner.get("Nama", "")
                    nama = str(nama_raw) if pd.notna(nama_raw) else ""
                    hp = format_phone(winner.get("No HP", ""))
                    display_nama = nama if nama and nama.lower() != "nan" else "-"
                    
                    st.markdown(f"""
                    <div style="background: linear-gradient(145deg, #fff, #f8f9fa); border-radius: 10px; padding: 0.6rem; text-align: center; border-left: 4px solid #f5576c; margin-bottom: 0.4rem; height: 75px; display: flex; flex-direction: column; justify-content: center;">
                        <div style="font-size: 1.1rem; font-weight: 800; color: #333; line-height: 1.3;">{nomor}</div>
                        <div style="font-size: 0.7rem; color: #666; line-height: 1.2; overflow: hidden; text-overflow: ellipsis; white-space: nowrap;">{display_nama}</div>
                        <div style="font-size: 0.65rem; color: #888; line-height: 1.2;">{hp}</div>
                    </div>
                    """, unsafe_allow_html=True)

elif current_page == "shuffle_page":
    remaining_pool = st.session_state.get("remaining_pool", pd.DataFrame())
    shuffle_results = st.session_state.get("shuffle_results", {})
    
    col_back, col_title, col_status = st.columns([1, 3, 2])
    with col_back:
        if st.button("⬅️ KEMBALI", key="back_from_shuffle"):
            st.session_state["current_page"] = "home"
            st.rerun()
    with col_title:
        st.markdown("<h2 style='text-align:center; color:#FF9800; margin:0;'>🎲 SHUFFLE</h2>", unsafe_allow_html=True)
    with col_status:
        done_count = len([k for k in shuffle_results.keys() if k.startswith("shuffle_batch")])
        st.markdown(f"<p style='text-align:right; color:#333; margin-top:10px;'>Sesi: <strong style='color:#FF9800;'>{done_count}/3</strong> | Sisa: <strong>{len(remaining_pool)}</strong></p>", unsafe_allow_html=True)
    
    shuffle_batches = [
        {"name": "Sesi 1", "count": 30},
        {"name": "Sesi 2", "count": 30},
        {"name": "Sesi 3", "count": 30},
    ]
    
    # Session tabs instead of expanders
    session_tabs = st.tabs(["📦 Sesi 1", "📦 Sesi 2", "📦 Sesi 3"])
    
    for i, batch in enumerate(shuffle_batches):
        batch_key = f"shuffle_batch_{i}"
        batch_done = batch_key in shuffle_results
        
        with session_tabs[i]:
            if batch_done:
                winners = shuffle_results[batch_key]["winners"]
                prize_assignments = shuffle_results[batch_key].get("prize_assignments", [])
                prize_config = shuffle_results[batch_key].get("prize_config", [])
                
                # Fallback for old format
                if not prize_assignments:
                    prize_name = shuffle_results[batch_key].get("prize_name", "Hadiah")
                    prize_assignments = [{"winner": w, "prize": prize_name} for w in winners]
                
                participant_data = st.session_state.get("participant_data")
                name_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["Nama"])) if participant_data is not None else {}
                phone_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["No HP"])) if participant_data is not None else {}
                
                # Create lookup for winner -> prize
                winner_prize_lookup = {pa["winner"]: pa["prize"] for pa in prize_assignments}
                
                # Group winners by prize
                prize_groups = {}
                for pa in prize_assignments:
                    prize = pa["prize"]
                    if prize not in prize_groups:
                        prize_groups[prize] = []
                    prize_groups[prize].append(pa["winner"])
                
                # Display each prize category
                for prize_name, prize_winners in prize_groups.items():
                    # Sort winners by nomor undian
                    sorted_winners = sorted(prize_winners, key=lambda x: str(x))
                    
                    # Header untuk kategori hadiah
                    st.markdown(f"""
                    <div style="background: linear-gradient(135deg, #4CAF50, #45a049); padding: 1rem; border-radius: 10px; text-align: center; margin: 1rem 0 0.5rem 0;">
                        <div style="font-size: 1.2rem; font-weight: bold; color: white;">🎁 {prize_name}</div>
                        <div style="font-size: 0.9rem; color: rgba(255,255,255,0.9);">{len(sorted_winners)} Pemenang</div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Display in 7 columns
                    num_cols = 7
                    rows = (len(sorted_winners) + num_cols - 1) // num_cols
                    
                    for row in range(rows):
                        row_cols = st.columns(num_cols)
                        for col in range(num_cols):
                            idx = row * num_cols + col
                            if idx < len(sorted_winners):
                                w = sorted_winners[idx]
                                with row_cols[col]:
                                    nama_raw = name_lookup.get(w, "")
                                    nama = str(nama_raw) if pd.notna(nama_raw) else ""
                                    display_nama = nama if nama and nama.lower() != "nan" else "-"
                                    hp = format_phone(phone_lookup.get(w, ""))
                                    st.markdown(f"""
                                    <div style="background: linear-gradient(145deg, #fff, #f8f9fa); border-radius: 10px; padding: 0.5rem; text-align: center; border-left: 4px solid #4CAF50; margin-bottom: 0.4rem; height: 70px; display: flex; flex-direction: column; justify-content: center;">
                                        <div style="font-size: 1rem; font-weight: 800; color: #333; line-height: 1.2;">{w}</div>
                                        <div style="font-size: 0.65rem; color: #666; line-height: 1.1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap;">{display_nama}</div>
                                        <div style="font-size: 0.6rem; color: #888; line-height: 1.1;">{hp}</div>
                                    </div>
                                    """, unsafe_allow_html=True)
                
                st.markdown("<br>", unsafe_allow_html=True)
                col1, col2 = st.columns(2)
                with col1:
                    # Build Excel with specific prize for each winner
                    excel_data = []
                    for pa in prize_assignments:
                        w = pa["winner"]
                        excel_data.append({
                            "Hadiah": pa["prize"],
                            "Nomor Undian": w,
                            "Nama": name_lookup.get(w, ""),
                            "No HP": phone_lookup.get(w, "")
                        })
                    df_batch = pd.DataFrame(excel_data)
                    # Sort by Hadiah then Nomor Undian
                    df_batch = df_batch.sort_values(["Hadiah", "Nomor Undian"])
                    excel_buf = BytesIO()
                    with pd.ExcelWriter(excel_buf, engine='openpyxl') as writer:
                        df_batch.to_excel(writer, index=False)
                    st.download_button(f"📊 Download Excel {batch['name']}", excel_buf.getvalue(), f"shuffle_{i+1}.xlsx", use_container_width=True)
                with col2:
                    pptx_data = generate_shuffle_pptx_v2(prize_assignments, name_lookup, phone_lookup, batch['name'])
                    st.download_button(f"📽️ Download PPT {batch['name']}", pptx_data, f"shuffle_{i+1}.pptx", use_container_width=True)
            else:
                remaining_count = len(remaining_pool)
                max_winners = min(batch['count'], remaining_count)
                
                st.markdown(f"**Konfigurasi Hadiah {batch['name']}** (Total: {max_winners} pemenang)")
                
                # Default shuffle prizes for this batch - different for each session
                shuffle_prize_key = f"shuffle_prizes_{batch_key}"
                if shuffle_prize_key not in st.session_state:
                    if i == 0:  # Sesi 1
                        st.session_state[shuffle_prize_key] = pd.DataFrame([
                            {"Nama Hadiah": "Sepeda Lipat (SJ-50MB-XB)", "Jumlah": 2},
                            {"Nama Hadiah": "Smart Watch Xiaomi (EO-35ST)", "Jumlah": 2},
                            {"Nama Hadiah": "Speaker (CBOX-B658UBO)", "Jumlah": 3},
                            {"Nama Hadiah": "Oven 18L (EO-18BL)", "Jumlah": 3},
                            {"Nama Hadiah": "Blender (EM-151G-GY)", "Jumlah": 4},
                            {"Nama Hadiah": "Rice Cooker (KS-N18MG-PK)", "Jumlah": 3},
                            {"Nama Hadiah": "Coffee Maker (HM-80L(W))", "Jumlah": 3},
                            {"Nama Hadiah": "Pop Up Toaster (KZ-2S02-BK)", "Jumlah": 4},
                            {"Nama Hadiah": "Hand Juicer (EM-P01-BK)", "Jumlah": 3},
                            {"Nama Hadiah": "Toaster (KZS-70L(W))", "Jumlah": 3},
                        ])
                    elif i == 1:  # Sesi 2
                        st.session_state[shuffle_prize_key] = pd.DataFrame([
                            {"Nama Hadiah": "Sepeda Lipat (SJ-50MB-XB)", "Jumlah": 2},
                            {"Nama Hadiah": "Smart Watch Xiaomi (EO-35ST)", "Jumlah": 1},
                            {"Nama Hadiah": "Speaker (CBOX-B658UBO)", "Jumlah": 3},
                            {"Nama Hadiah": "Oven 18L (EO-18BL)", "Jumlah": 4},
                            {"Nama Hadiah": "Blender (EM-151G-GY)", "Jumlah": 3},
                            {"Nama Hadiah": "Rice Cooker (KS-N18MG-PK)", "Jumlah": 3},
                            {"Nama Hadiah": "Coffee Maker (HM-80L(W))", "Jumlah": 4},
                            {"Nama Hadiah": "Pop Up Toaster (KZ-2S02-BK)", "Jumlah": 3},
                            {"Nama Hadiah": "Hand Juicer (EM-P01-BK)", "Jumlah": 3},
                            {"Nama Hadiah": "Toaster (KZS-70L(W))", "Jumlah": 4},
                        ])
                    else:  # Sesi 3
                        st.session_state[shuffle_prize_key] = pd.DataFrame([
                            {"Nama Hadiah": "Sepeda Lipat (SJ-50MB-XB)", "Jumlah": 1},
                            {"Nama Hadiah": "Smart Watch Xiaomi (EO-35ST)", "Jumlah": 2},
                            {"Nama Hadiah": "Speaker (CBOX-B658UBO)", "Jumlah": 4},
                            {"Nama Hadiah": "Oven 18L (EO-18BL)", "Jumlah": 3},
                            {"Nama Hadiah": "Blender (EM-151G-GY)", "Jumlah": 3},
                            {"Nama Hadiah": "Rice Cooker (KS-N18MG-PK)", "Jumlah": 4},
                            {"Nama Hadiah": "Coffee Maker (HM-80L(W))", "Jumlah": 3},
                            {"Nama Hadiah": "Pop Up Toaster (KZ-2S02-BK)", "Jumlah": 3},
                            {"Nama Hadiah": "Hand Juicer (EM-P01-BK)", "Jumlah": 4},
                            {"Nama Hadiah": "Toaster (KZS-70L(W))", "Jumlah": 3},
                        ])
                
                edited_prizes = st.data_editor(
                    st.session_state[shuffle_prize_key],
                    num_rows="dynamic",
                    use_container_width=True,
                    key=f"editor_{batch_key}",
                    column_config={
                        "Nama Hadiah": st.column_config.TextColumn("Nama Hadiah", width="large"),
                        "Jumlah": st.column_config.NumberColumn("Jumlah", min_value=1, max_value=100, width="small")
                    }
                )
                st.session_state[shuffle_prize_key] = edited_prizes
                
                total_prizes = edited_prizes["Jumlah"].sum() if len(edited_prizes) > 0 else 0
                
                if total_prizes != max_winners:
                    st.warning(f"⚠️ Total hadiah ({int(total_prizes)}) harus sama dengan jumlah pemenang ({max_winners})")
                else:
                    st.success(f"✅ Total hadiah: {int(total_prizes)} = {max_winners} pemenang")
                
                if remaining_count > 0 and total_prizes == max_winners and len(edited_prizes) > 0:
                    if st.button(f"🎲 MULAI {batch['name']}", key=f"start_{batch_key}", use_container_width=True):
                        remaining_numbers = remaining_pool["Nomor Undian"].tolist()
                        batch_winners = []
                        temp_pool = remaining_numbers.copy()
                        
                        for _ in range(max_winners):
                            if len(temp_pool) == 0:
                                break
                            idx = secrets.randbelow(len(temp_pool))
                            batch_winners.append(temp_pool.pop(idx))
                        
                        # Show shuffle animation
                        scroll_js = """
                        <script>
                            setTimeout(function() {
                                var animElement = document.querySelector('iframe');
                                if (animElement) {
                                    animElement.scrollIntoView({behavior: 'smooth', block: 'center'});
                                }
                            }, 100);
                        </script>
                        """
                        components.html(scroll_js, height=0)
                        
                        shuffle_html = create_shuffle_animation_html(remaining_numbers, batch_winners, batch['name'])
                        components.html(shuffle_html, height=420)
                        
                        # Assign prizes to winners
                        prize_assignments = []
                        winner_idx = 0
                        for _, row in edited_prizes.iterrows():
                            prize_name = row["Nama Hadiah"]
                            count = int(row["Jumlah"])
                            for _ in range(count):
                                if winner_idx < len(batch_winners):
                                    prize_assignments.append({
                                        "winner": batch_winners[winner_idx],
                                        "prize": prize_name
                                    })
                                    winner_idx += 1
                        
                        shuffle_results[batch_key] = {
                            "winners": batch_winners,
                            "prize_assignments": prize_assignments,
                            "prize_config": edited_prizes.to_dict('records')
                        }
                        st.session_state["shuffle_results"] = shuffle_results
                        
                        new_pool = remaining_pool[~remaining_pool["Nomor Undian"].isin(batch_winners)]
                        st.session_state["remaining_pool"] = new_pool
                        
                        if len(shuffle_results) == 3:
                            st.session_state["shuffle_done"] = True
                        
                        # Auto-save results
                        save_lottery_results()
                        
                        # Show success message after animation
                        time.sleep(0.5)
                        st.success(f"🎉 {len(batch_winners)} pemenang {batch['name']} berhasil diundi!")
                        
                        if st.button("✅ Lihat Hasil Lengkap", key=f"view_result_{batch_key}", use_container_width=True):
                            st.rerun()
                elif remaining_count == 0:
                    st.warning("Tidak ada sisa peserta")
    
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("---")
    
    with st.expander(f"📋 Nomor yang Belum Diundi ({len(remaining_pool)} peserta)", expanded=False):
        if len(remaining_pool) > 0:
            remaining_numbers = remaining_pool["Nomor Undian"].tolist()
            cols = st.columns(15)
            for idx, num in enumerate(remaining_numbers[:150]):
                with cols[idx % 15]:
                    st.markdown(f"<div style='background:#333;color:white;padding:0.3rem;border-radius:5px;text-align:center;margin:2px;font-size:0.8rem;'>{num}</div>", unsafe_allow_html=True)
            if len(remaining_numbers) > 150:
                st.info(f"... dan {len(remaining_numbers) - 150} nomor lainnya")
        else:
            st.info("Semua nomor sudah diundi")
    
    if st.button("📊 SISA NOMOR → KEMBALI KE MENU UTAMA", key="shuffle_done_btn", use_container_width=True):
        st.session_state["current_page"] = "home"
        st.rerun()

elif current_page == "shuffle_results":
    batch_idx = st.session_state.get("viewing_shuffle_batch", 0)
    shuffle_results = st.session_state.get("shuffle_results", {})
    batch_key = f"shuffle_batch_{batch_idx}"
    
    if batch_key not in shuffle_results:
        st.session_state["current_page"] = "home"
        st.rerun()
    
    if st.button("⬅️ KEMBALI", key="back_from_shuffle_results"):
        st.session_state["current_page"] = "home"
        st.rerun()
    
    result = shuffle_results[batch_key]
    winners = result["winners"]
    prize_name = result["prize_name"]
    
    st.markdown(f"""
    <div style="background: linear-gradient(135deg, #FF9800, #FF5722); padding: 2rem; border-radius: 15px; text-align: center; margin: 1rem 0;">
        <p style="color: white; font-size: 2.5rem; font-weight: bold; margin: 0;">🎲 {prize_name}</p>
        <p style="color: #fff; font-size: 1.2rem; margin: 0.5rem 0;">Sesi {batch_idx + 1} - {len(winners)} Pemenang</p>
    </div>
    """, unsafe_allow_html=True)
    
    participant_data = st.session_state.get("participant_data")
    name_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["Nama"])) if participant_data is not None else {}
    phone_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["No HP"])) if participant_data is not None else {}
    
    cols = 10
    rows = (len(winners) + cols - 1) // cols
    
    for row in range(rows):
        row_cols = st.columns(cols)
        for col in range(cols):
            idx = row * cols + col
            if idx < len(winners):
                w = winners[idx]
                with row_cols[col]:
                    nama_raw = name_lookup.get(w, "")
                    nama = str(nama_raw) if pd.notna(nama_raw) else ""
                    display_nama = nama if nama and nama.lower() != "nan" else "-"
                    hp = format_phone(phone_lookup.get(w, ""))
                    st.markdown(f"""
                    <div style="background: linear-gradient(145deg, #fff, #f8f9fa); border-radius: 10px; padding: 0.8rem; text-align: center; border-left: 4px solid #FF9800; margin-bottom: 0.5rem;">
                        <div style="font-size: 1.2rem; font-weight: 800; color: #333;">{w}</div>
                        <div style="font-size: 0.7rem; color: #666;">{display_nama}</div>
                        <div style="font-size: 0.7rem; color: #888;">{hp}</div>
                    </div>
                    """, unsafe_allow_html=True)

elif current_page == "wheel_page":
    remaining_pool = st.session_state.get("remaining_pool", pd.DataFrame())
    wheel_winners = st.session_state.get("wheel_winners", [])
    wheel_prizes = st.session_state.get("wheel_prizes", [])
    
    col_back, col_title, col_status = st.columns([1, 3, 2])
    with col_back:
        if st.button("⬅️ KEMBALI", key="back_from_wheel"):
            st.session_state["current_page"] = "home"
            st.rerun()
    with col_title:
        st.markdown("<h2 style='text-align:center; color:#E91E63; margin:0;'>🎡 HADIAH UTAMA</h2>", unsafe_allow_html=True)
    with col_status:
        st.markdown(f"<p style='text-align:right; color:#333; margin-top:10px;'>Pemenang: <strong style='color:#E91E63;'>{len(wheel_winners)}/10</strong></p>", unsafe_allow_html=True)
    
    def get_valid_wheel_config():
        """Ensure wheel_config has the correct format"""
        default_config = [
            {"No": 1, "Nama Hadiah": "HP Samsung A07", "Keterangan": "EC-8305-B"},
            {"No": 2, "Nama Hadiah": "HP Samsung A07", "Keterangan": "EC-8305-B"},
            {"No": 3, "Nama Hadiah": "Kulkas 1 Pintu", "Keterangan": "SJ-N162D-AP"},
            {"No": 4, "Nama Hadiah": "Kulkas 1 Pintu", "Keterangan": "SJ-N162D-AP"},
            {"No": 5, "Nama Hadiah": "Mesin Cuci Matic 7KG", "Keterangan": "ES-M7000P-GG"},
            {"No": 6, "Nama Hadiah": "Mesin Cuci Matic 7KG", "Keterangan": "ES-M7000P-GG"},
            {"No": 7, "Nama Hadiah": "Mesin Cuci Matic 7KG", "Keterangan": "ES-M7000P-GG"},
            {"No": 8, "Nama Hadiah": "Mesin Cuci Matic 7KG", "Keterangan": "ES-M7000P-GG"},
            {"No": 9, "Nama Hadiah": "LED TV 43\"", "Keterangan": "43HJ6000I"},
            {"No": 10, "Nama Hadiah": "LED TV 43\"", "Keterangan": "43HJ6000I"},
        ]
        
        existing = st.session_state.get("wheel_config", [])
        
        if not existing or len(existing) == 0:
            return default_config
        
        if "Nama Hadiah" not in existing[0]:
            new_config = []
            for i, item in enumerate(existing[:10]):
                prize_name = item.get("prize", "") or item.get("name", "") or f"Grand Prize {i+1}"
                new_config.append({
                    "No": i+1,
                    "Nama Hadiah": prize_name if prize_name else f"Grand Prize {i+1}",
                    "Keterangan": ""
                })
            while len(new_config) < 10:
                new_config.append({
                    "No": len(new_config)+1,
                    "Nama Hadiah": f"Grand Prize {len(new_config)+1}",
                    "Keterangan": ""
                })
            return new_config
        
        return existing
    
    st.session_state["wheel_config"] = get_valid_wheel_config()
    
    with st.expander("⚙️ Edit Hadiah", expanded=False):
        wheel_config_df = pd.DataFrame(st.session_state["wheel_config"])
        edited_wheel_config = st.data_editor(
            wheel_config_df, num_rows="fixed", use_container_width=True, hide_index=True,
            column_config={
                "No": st.column_config.NumberColumn("No", disabled=True, width="small"),
                "Nama Hadiah": st.column_config.TextColumn("Nama Hadiah", width="medium", required=True),
                "Keterangan": st.column_config.TextColumn("Keterangan", width="large")
            },
            key="wheel_config_editor"
        )
        st.session_state["wheel_config"] = edited_wheel_config.to_dict('records')
    
    current_idx = len(wheel_winners)
    wheel_config = st.session_state.get("wheel_config", [])
    is_spinning = st.session_state.get("wheel_spinning", False)
    
    # Compact prize progress bar
    progress_html = "<div style='display:flex; gap:3px; justify-content:center; margin:5px 0;'>"
    for i in range(10):
        if i < len(wheel_winners):
            progress_html += f"<span style='background:#4CAF50;color:white;padding:3px 8px;border-radius:3px;font-size:0.7rem;'>✓{i+1}</span>"
        elif i == current_idx:
            progress_html += f"<span style='background:#E91E63;color:white;padding:3px 8px;border-radius:3px;font-size:0.7rem;font-weight:bold;'>▶{i+1}</span>"
        else:
            progress_html += f"<span style='background:#ddd;color:#999;padding:3px 8px;border-radius:3px;font-size:0.7rem;'>{i+1}</span>"
    progress_html += "</div>"
    st.markdown(progress_html, unsafe_allow_html=True)
    
    if current_idx < 10 and len(remaining_pool) > 0:
        prize_name = wheel_config[current_idx].get("Nama Hadiah", f"Prize {current_idx + 1}") if current_idx < len(wheel_config) else f"Prize {current_idx + 1}"
        prize_keterangan = wheel_config[current_idx].get("Keterangan", "") if current_idx < len(wheel_config) else ""
        
        # Prize info + button in one compact row
        st.markdown(f"""
        <div style="background:linear-gradient(135deg,#E91E63,#9C27B0);border-radius:10px;padding:10px 15px;text-align:center;margin:5px 0;">
            <span style="color:white;font-size:1.2rem;font-weight:bold;">🎁 #{current_idx+1}: {prize_name}</span>
            <span style="color:rgba(255,255,255,0.8);font-size:0.85rem;margin-left:10px;">{prize_keterangan}</span>
        </div>
        """, unsafe_allow_html=True)
        
        if st.button("🎡 PUTAR UNDIAN!", key=f"spin_wheel_{current_idx}", use_container_width=True):
            remaining_numbers = remaining_pool["Nomor Undian"].tolist()
            
            if len(remaining_numbers) > 0:
                winner_idx = secrets.randbelow(len(remaining_numbers))
                winner = remaining_numbers[winner_idx]
                
                wheel_html = create_spinning_wheel_html(remaining_numbers, winner, 320)
                components.html(wheel_html, height=280)
                
                wheel_winners.append(winner)
                wheel_prizes.append(prize_name)
                st.session_state["wheel_winners"] = wheel_winners
                st.session_state["wheel_prizes"] = wheel_prizes
                
                new_pool = remaining_pool[remaining_pool["Nomor Undian"] != winner]
                st.session_state["remaining_pool"] = new_pool
                
                if len(wheel_winners) == 10:
                    st.session_state["wheel_done"] = True
                
                save_lottery_results()
                
                participant_data = st.session_state.get("participant_data")
                if participant_data is not None:
                    winner_row = participant_data[participant_data["Nomor Undian"] == winner]
                    if len(winner_row) > 0:
                        nama_raw = winner_row.iloc[0].get("Nama", "")
                        nama = str(nama_raw) if pd.notna(nama_raw) else "-"
                        if nama.lower() == "nan":
                            nama = "-"
                        hp = format_phone(winner_row.iloc[0].get("No HP", ""))
                        st.markdown(f"""
                        <div style="background:#4CAF50;color:white;padding:15px;border-radius:10px;text-align:center;margin:10px 0;">
                            <div style="font-size:1.5rem;font-weight:bold;">🎉 PEMENANG #{current_idx+1}</div>
                            <div style="font-size:2rem;font-weight:900;margin:5px 0;">{winner}</div>
                            <div style="font-size:1rem;">{nama} | {hp}</div>
                        </div>
                        """, unsafe_allow_html=True)
                
                if len(wheel_winners) < 10:
                    st.button("➡️ LANJUT KE HADIAH BERIKUTNYA", key="next_wheel", use_container_width=True)
    
    # Previous winners - full width cards
    if len(wheel_winners) > 0:
        st.markdown("---")
        participant_data = st.session_state.get("participant_data")
        name_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["Nama"])) if participant_data is not None else {}
        phone_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["No HP"])) if participant_data is not None else {}
        
        # Full width winner cards in 5 columns
        cols = st.columns(5)
        for i, (w, p) in enumerate(zip(wheel_winners, wheel_prizes)):
            nama_raw = name_lookup.get(w, "")
            nama = str(nama_raw) if pd.notna(nama_raw) else "-"
            if nama.lower() == "nan":
                nama = "-"
            hp = format_phone(phone_lookup.get(w, ""))
            with cols[i % 5]:
                st.markdown(f"""
                <div style="background:#fff;border:2px solid #E91E63;border-radius:10px;padding:8px;text-align:center;margin-bottom:5px;">
                    <div style="font-size:0.75rem;color:#E91E63;font-weight:bold;">#{i+1}</div>
                    <div style="font-size:1.1rem;font-weight:800;color:#333;">{w}</div>
                    <div style="font-size:0.7rem;color:#666;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;">{nama}</div>
                    <div style="font-size:0.65rem;color:#888;">{hp}</div>
                    <div style="font-size:0.6rem;color:#E91E63;margin-top:3px;">{p}</div>
                </div>
                """, unsafe_allow_html=True)
        
        # Download buttons
        col1, col2 = st.columns(2)
        with col1:
            df_wheel = pd.DataFrame({
                "No": range(1, len(wheel_winners) + 1),
                "Nomor Undian": wheel_winners,
                "Nama": [name_lookup.get(w, "") for w in wheel_winners],
                "No HP": [phone_lookup.get(w, "") for w in wheel_winners],
                "Hadiah": wheel_prizes
            })
            excel_buf = BytesIO()
            with pd.ExcelWriter(excel_buf, engine='openpyxl') as writer:
                df_wheel.to_excel(writer, index=False)
            st.download_button("📊 Download Excel Wheel", excel_buf.getvalue(), "wheel_winners.xlsx", use_container_width=True)
        
        with col2:
            pptx_data = generate_wheel_pptx(wheel_winners, wheel_prizes, name_lookup, phone_lookup)
            st.download_button("📽️ Download PPT Wheel", pptx_data, "wheel_winners.pptx", use_container_width=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        with st.expander(f"📋 Nomor yang Belum Diundi ({len(remaining_pool)} peserta)", expanded=False):
            if len(remaining_pool) > 0:
                remaining_numbers = remaining_pool["Nomor Undian"].tolist()
                cols = st.columns(15)
                for idx, num in enumerate(remaining_numbers[:150]):
                    with cols[idx % 15]:
                        st.markdown(f"<div style='background:#333;color:white;padding:0.3rem;border-radius:5px;text-align:center;margin:2px;font-size:0.8rem;'>{num}</div>", unsafe_allow_html=True)
                if len(remaining_numbers) > 150:
                    st.info(f"... dan {len(remaining_numbers) - 150} nomor lainnya")
            else:
                st.info("Semua nomor sudah diundi")
        
        # Final buttons (only show when wheel is complete)
        if st.session_state.get("wheel_done", False):
            st.markdown("---")
            
            # Compact 3-button row
            val_col, excel_col, ppt_col = st.columns(3)
            
            with val_col:
                if st.button("🔍 CEK VALIDASI", key="validate_all", use_container_width=True):
                    all_winners = []
                    duplicate_info = []
                    
                    evoucher_results = st.session_state.get("evoucher_results")
                    if evoucher_results is not None and len(evoucher_results) > 0:
                        for num in evoucher_results["Nomor Undian"].tolist():
                            if num in all_winners:
                                duplicate_info.append(f"E-Voucher: {num}")
                            all_winners.append(num)
                    
                    shuffle_results = st.session_state.get("shuffle_results", {})
                    for batch_key, batch_data in shuffle_results.items():
                        for num in batch_data.get("winners", []):
                            if num in all_winners:
                                duplicate_info.append(f"Shuffle: {num}")
                            all_winners.append(num)
                    
                    for num in wheel_winners:
                        if num in all_winners:
                            duplicate_info.append(f"Wheel: {num}")
                        all_winners.append(num)
                    
                    if len(duplicate_info) > 0:
                        st.error(f"⚠️ {len(duplicate_info)} DOBEL!")
                    else:
                        st.success(f"✅ OK! {len(all_winners)} pemenang unik")
                        st.balloons()
            
            with excel_col:
                combined_excel = BytesIO()
                with pd.ExcelWriter(combined_excel, engine='openpyxl') as writer:
                    evoucher_results = st.session_state.get("evoucher_results")
                    if evoucher_results is not None and len(evoucher_results) > 0:
                        evoucher_results.to_excel(writer, sheet_name="E-Voucher", index=False)
                    
                    shuffle_results = st.session_state.get("shuffle_results", {})
                    for batch_key, batch_data in shuffle_results.items():
                        batch_winners = batch_data.get("winners", [])
                        if len(batch_winners) > 0:
                            df_batch = pd.DataFrame({
                                "No": range(1, len(batch_winners) + 1),
                                "Nomor Undian": batch_winners,
                                "Nama": [name_lookup.get(w, "") for w in batch_winners],
                                "No HP": [phone_lookup.get(w, "") for w in batch_winners],
                                "Hadiah": [batch_data.get("prize_name", "")] * len(batch_winners)
                            })
                            df_batch.to_excel(writer, sheet_name=f"Shuffle_{batch_key.split('_')[-1]}", index=False)
                    
                    if len(wheel_winners) > 0:
                        pd.DataFrame({
                            "No": range(1, len(wheel_winners) + 1),
                            "Nomor Undian": wheel_winners,
                            "Nama": [name_lookup.get(w, "") for w in wheel_winners],
                            "No HP": [phone_lookup.get(w, "") for w in wheel_winners],
                            "Hadiah": wheel_prizes
                        }).to_excel(writer, sheet_name="Grand_Prize", index=False)
                
                st.download_button("📊 EXCEL LENGKAP", combined_excel.getvalue(), "MoveGroove_Lengkap.xlsx", use_container_width=True)
            
            with ppt_col:
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                slide_layout = prs.slide_layouts[6]
                
                # Title slide
                slide = prs.slides.add_slide(slide_layout)
                shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(33, 150, 243)
                shape.line.fill.background()
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(2))
                tf = title_box.text_frame
                tf.paragraphs[0].text = "MOVE & GROOVE 2024"
                tf.paragraphs[0].font.size = Pt(60)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                p2 = tf.add_paragraph()
                p2.text = "HASIL UNDIAN LENGKAP"
                p2.font.size = Pt(36)
                p2.font.color.rgb = RGBColor(255, 255, 255)
                p2.alignment = PP_ALIGN.CENTER
                
                # E-Voucher slides
                evoucher_results = st.session_state.get("evoucher_results")
                if evoucher_results is not None and len(evoucher_results) > 0:
                    for tier in st.session_state.get("prize_tiers", []):
                        tier_winners = evoucher_results[evoucher_results["Kategori"] == tier["name"]]["Nomor Undian"].tolist()
                        if tier_winners:
                            slide = prs.slides.add_slide(slide_layout)
                            shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(76, 175, 80)
                            shape.line.fill.background()
                            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
                            title_box.text_frame.paragraphs[0].text = f"E-VOUCHER: {tier['name']}"
                            title_box.text_frame.paragraphs[0].font.size = Pt(32)
                            title_box.text_frame.paragraphs[0].font.bold = True
                            title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                            title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            for idx, w in enumerate(tier_winners):
                                row, col = idx // 10, idx % 10
                                cell = slide.shapes.add_shape(5, Inches(0.5) + col * Inches(1.28), Inches(1.3) + row * Inches(0.58), Inches(1.2), Inches(0.5))
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                                cell.text_frame.paragraphs[0].text = str(w)
                                cell.text_frame.paragraphs[0].font.size = Pt(14)
                                cell.text_frame.paragraphs[0].font.bold = True
                                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Shuffle slides
                for batch_key, batch_data in st.session_state.get("shuffle_results", {}).items():
                    batch_winners = batch_data.get("winners", [])
                    if batch_winners:
                        slide = prs.slides.add_slide(slide_layout)
                        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(156, 39, 176)
                        shape.line.fill.background()
                        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
                        title_box.text_frame.paragraphs[0].text = f"SHUFFLE: {batch_data.get('prize_name', '')}"
                        title_box.text_frame.paragraphs[0].font.size = Pt(32)
                        title_box.text_frame.paragraphs[0].font.bold = True
                        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        for idx, w in enumerate(batch_winners):
                            row, col = idx // 10, idx % 10
                            cell = slide.shapes.add_shape(5, Inches(0.5) + col * Inches(1.28), Inches(1.3) + row * Inches(0.58), Inches(1.2), Inches(0.5))
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            cell.text_frame.paragraphs[0].text = str(w)
                            cell.text_frame.paragraphs[0].font.size = Pt(14)
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Wheel slide
                if wheel_winners:
                    slide = prs.slides.add_slide(slide_layout)
                    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(233, 30, 99)
                    shape.line.fill.background()
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
                    title_box.text_frame.paragraphs[0].text = "GRAND PRIZE - SPINNING WHEEL"
                    title_box.text_frame.paragraphs[0].font.size = Pt(32)
                    title_box.text_frame.paragraphs[0].font.bold = True
                    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    for idx, (w, p) in enumerate(zip(wheel_winners, wheel_prizes)):
                        cell = slide.shapes.add_shape(5, Inches(1) + (idx % 2) * Inches(6), Inches(1.5) + (idx // 2) * Inches(1.1), Inches(5.5), Inches(1))
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        nama = str(name_lookup.get(w, "")) if pd.notna(name_lookup.get(w, "")) else "-"
                        cell.text_frame.paragraphs[0].text = f"#{idx+1} {w} - {nama[:25]}"
                        cell.text_frame.paragraphs[0].font.size = Pt(18)
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 30, 99)
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        p2 = cell.text_frame.add_paragraph()
                        p2.text = p[:40]
                        p2.font.size = Pt(14)
                        p2.font.color.rgb = RGBColor(100, 100, 100)
                        p2.alignment = PP_ALIGN.CENTER
                
                ppt_buffer = BytesIO()
                prs.save(ppt_buffer)
                st.download_button("📽️ PPT LENGKAP", ppt_buffer.getvalue(), "MoveGroove_Lengkap.pptx", use_container_width=True)
        
        if st.button("🏠 KEMBALI KE MENU UTAMA", key="wheel_done_btn", use_container_width=True):
            st.session_state["current_page"] = "home"
            st.rerun()

elif current_page == "wheel_results":
    wheel_winners = st.session_state.get("wheel_winners", [])
    wheel_prizes = st.session_state.get("wheel_prizes", [])
    
    if len(wheel_winners) == 0:
        st.session_state["current_page"] = "home"
        st.rerun()
    
    if st.button("⬅️ KEMBALI", key="back_from_wheel_results"):
        st.session_state["current_page"] = "home"
        st.rerun()
    
    st.markdown("""
    <div style="background: linear-gradient(135deg, #E91E63, #9C27B0); padding: 2rem; border-radius: 15px; text-align: center; margin: 1rem 0;">
        <p style="color: white; font-size: 2.5rem; font-weight: bold; margin: 0;">🎡 PEMENANG SPINNING WHEEL</p>
        <p style="color: #fff; font-size: 1.2rem; margin: 0.5rem 0;">{len(wheel_winners)} Grand Prize Winners</p>
    </div>
    """, unsafe_allow_html=True)
    
    participant_data = st.session_state.get("participant_data")
    name_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["Nama"])) if participant_data is not None else {}
    phone_lookup = dict(zip(participant_data["Nomor Undian"], participant_data["No HP"])) if participant_data is not None else {}
    
    cols = st.columns(5)
    for i, (w, p) in enumerate(zip(wheel_winners, wheel_prizes)):
        with cols[i % 5]:
            nama_raw = name_lookup.get(w, "")
            nama = str(nama_raw) if pd.notna(nama_raw) else ""
            display_nama = nama if nama and nama.lower() != "nan" else "-"
            hp = format_phone(phone_lookup.get(w, ""))
            st.markdown(f"""
            <div style="background: linear-gradient(145deg, #fff, #f8f9fa); border-radius: 15px; padding: 1rem; text-align: center; border: 3px solid #E91E63; margin-bottom: 1rem;">
                <div style="font-size: 0.9rem; color: #E91E63; font-weight: bold;">#{i+1}</div>
                <div style="font-size: 1.5rem; font-weight: 800; color: #333;">{w}</div>
                <div style="font-size: 0.85rem; color: #666;">{display_nama}</div>
                <div style="font-size: 0.8rem; color: #888;">{hp}</div>
                <div style="font-size: 0.75rem; color: #E91E63; margin-top: 0.5rem;">{p}</div>
            </div>
            """, unsafe_allow_html=True)
